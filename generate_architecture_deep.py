"""
Generate Oil&Gas Downstream POV - Deep Architecture Section
Enriching the architecture narrative with technical depth
"""
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- CONSTANTS ---
SLIDE_WIDTH = Emu(18288000)
SLIDE_HEIGHT = Emu(10287000)

COLOR_DK1 = RGBColor(0x00, 0x1E, 0x5B)
COLOR_DK2 = RGBColor(0x02, 0x2A, 0xC0)
COLOR_ACCENT1 = RGBColor(0x90, 0xD0, 0xFE)
COLOR_ACCENT2 = RGBColor(0x04, 0xE1, 0xCB)
COLOR_ACCENT6 = RGBColor(0xEA, 0xF5, 0xFE)
COLOR_HEADER = RGBColor(0x12, 0x37, 0x7A)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_MUTED = RGBColor(0x85, 0x84, 0x81)
COLOR_CARD_BODY = RGBColor(0x38, 0x46, 0x53)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_BG_LIGHT = RGBColor(0xCF, 0xE9, 0xFE)
COLOR_GREEN = RGBColor(0x2E, 0x7D, 0x32)
COLOR_PURPLE = RGBColor(0x6A, 0x1B, 0x9A)
COLOR_ORANGE = RGBColor(0xE6, 0x5C, 0x00)

FONT_TITLE = "Avant Garde Demi SFDC"
FONT_BODY = "Salesforce Sans"
FONT_CARD = "Inter"
FONT_CONTENT = "Montserrat"
MARGIN_LEFT = Emu(900000)


def add_text_block(slide, left, top, width, height, text, font_name=FONT_BODY,
                   font_size=Pt(16), color=COLOR_BODY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return tf


def add_multi_text(slide, left, top, width, height, lines, font_name=FONT_BODY,
                   font_size=Pt(16), color=COLOR_BODY, spacing=Pt(8)):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color
        p.space_after = spacing
    return tf


def add_section_divider(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_DK1

    txBox = slide.shapes.add_textbox(Emu(1800000), Emu(3200000), Emu(14400000), Emu(2400000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(44)
    run.font.color.rgb = COLOR_WHITE

    txBox2 = slide.shapes.add_textbox(Emu(1800000), Emu(5800000), Emu(14400000), Emu(1600000))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_BODY
    run2.font.size = Pt(22)
    run2.font.color.rgb = COLOR_ACCENT1
    return slide


def add_content_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = COLOR_WHITE

    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Emu(500000), Emu(16000000), Emu(1200000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(34)
    run.font.color.rgb = COLOR_DK1

    if subtitle:
        txBox2 = slide.shapes.add_textbox(MARGIN_LEFT, Emu(1300000), Emu(16000000), Emu(900000))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_BODY
        run2.font.size = Pt(18)
        run2.font.color.rgb = COLOR_BODY
    return slide


def add_footer(slide, text="Salesforce POV • Architecture Deep Dive"):
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Emu(9700000), Emu(8000000), Emu(400000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CARD
    run.font.size = Pt(11)
    run.font.color.rgb = COLOR_MUTED


def add_layer_bar(slide, left, top, width, height, text, color, subtext=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_block(slide, left + Emu(200000), top + Emu(80000), width - Emu(400000), Emu(400000),
                  text, FONT_BODY, Pt(14), COLOR_WHITE, bold=True)
    if subtext:
        add_text_block(slide, left + Emu(200000), top + Emu(420000), width - Emu(400000), Emu(400000),
                      subtext, FONT_CONTENT, Pt(12), RGBColor(0xCC, 0xE0, 0xFF))


# ============================================================
# BUILD
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT


# ============================================================
# SLIDE 1: ARCHITECTURE SECTION DIVIDER
# ============================================================
add_section_divider(prs,
    "Platform Architecture Deep Dive",
    "How Salesforce, MuleSoft, Data Cloud, Tableau and Agentforce compose the downstream intelligence layer"
)


# ============================================================
# SLIDE 2: REFERENCE ARCHITECTURE (Full Stack)
# ============================================================
slide = add_content_slide(prs,
    "Target Reference Architecture",
    "Five architectural layers connect systems of record to governed commercial and operational action"
)

layers = [
    (COLOR_PURPLE, "AGENTIC LAYER",
     "Agentforce — Autonomous agents for pricing, logistics, commercial, credit and retention decisions",
     "Monitors triggers • Reasons over context • Recommends action • Documents governance • Learns from outcomes"),
    (COLOR_DK2, "ENGAGEMENT & EXECUTION LAYER",
     "Sales Cloud • Revenue Cloud • Service Cloud • Field Service • Experience Cloud • Loyalty Management",
     "Account plans • Quotes • Cases • Work orders • Dealer portals • Fleet portals • Consumer loyalty • Approvals"),
    (COLOR_HEADER, "ANALYTICS & INTELLIGENCE LAYER",
     "Tableau • CRM Analytics • Einstein Predictions",
     "Executive dashboards • Operational KPIs • Churn models • Margin analysis • Demand forecasting • Anomaly detection"),
    (COLOR_DK1, "DATA UNIFICATION LAYER",
     "Data Cloud (Customer 360) — Unified profiles for accounts, sites, assets, products, contracts and interactions",
     "Identity resolution • Segmentation • Calculated insights • Real-time streaming • Activation to engagement layer"),
    (COLOR_GREEN, "INTEGRATION LAYER",
     "MuleSoft Anypoint Platform — API-led connectivity and event-driven architecture",
     "System APIs (SAP, TMS, Fiscal) • Process APIs (pricing, logistics, credit) • Experience APIs (portals, mobile, agents)"),
]

y = Emu(2400000)
for i, (color, name, products, capabilities) in enumerate(layers):
    row_y = y + i * Emu(1450000)
    # Layer bar
    add_layer_bar(slide, Emu(900000), row_y, Emu(16400000), Emu(1250000), name, color)
    # Products
    add_text_block(slide, Emu(4000000), row_y + Emu(80000), Emu(13000000), Emu(400000),
                  products, FONT_CARD, Pt(13), COLOR_WHITE, bold=False)
    # Capabilities
    add_text_block(slide, Emu(4000000), row_y + Emu(500000), Emu(13000000), Emu(600000),
                  capabilities, FONT_CARD, Pt(11), RGBColor(0xCC, 0xE0, 0xFF))

# Source systems at bottom
add_text_block(slide, Emu(900000), Emu(9600000), Emu(16400000), Emu(400000),
              "SOURCE SYSTEMS:  SAP IS-Oil&Gas  •  SAP S/4 (MM, FI, CO-PA, SD)  •  TMS  •  IoT/Telemetry  •  SEFAZ/NF-e  •  ANP  •  Petrobras Feeds  •  Carriers EDI  •  Legacy Loyalty",
              FONT_CARD, Pt(11), COLOR_MUTED)


# ============================================================
# SLIDE 3: INDUSTRY CLOUD DECISION
# ============================================================
slide = add_content_slide(prs,
    "Industry Cloud Strategy: Composable, Not Monolithic",
    "Downstream fuel distribution combines patterns from Energy, Manufacturing and Consumer Goods — the architecture should too"
)

add_text_block(slide, MARGIN_LEFT, Emu(2300000), Emu(16000000), Emu(700000),
              "Why no single Industry Cloud is a perfect fit for downstream fuel distribution:",
              FONT_BODY, Pt(18), COLOR_DK1, bold=True)

# Three columns - one per industry cloud
clouds = [
    ("Manufacturing Cloud", COLOR_DK2,
     "STRONGEST FIT FOR:",
     ["Sales Agreements (B2B volume contracts)", "Account-Based Forecasting (demand by client)",
      "Rebate Management (dealer incentive programs)", "Run-rate business visibility"],
     "USE WHEN:",
     "Managing B2B contracts with agribusiness, mining, fleets. Forecasting demand. Running dealer incentive/rebate programs."),
    ("Energy & Utilities Cloud", COLOR_GREEN,
     "STRONGEST FIT FOR:",
     ["Account-to-Location data model (stations/sites)", "Asset management (tanks, terminals, PAVs)",
      "Consumption/usage tracking patterns", "Regulated operations & compliance"],
     "USE WHEN:",
     "Managing physical assets (terminals, tanks). Tracking consumption by site. Operating in regulated context (ANP, CBIOs)."),
    ("Consumer Goods Cloud", COLOR_PURPLE,
     "STRONGEST FIT FOR:",
     ["Retail Execution (structured dealer visits)", "Trade Promotion Management",
      "Key Account Management for retail networks", "Field activity and compliance audits"],
     "USE WHEN:",
     "Managing 44K station visits. Running promotions. Auditing brand compliance. Managing convenience retail (AmPm, Shell Select).")
]

for i, (cloud_name, color, fit_label, fit_items, use_label, use_text) in enumerate(clouds):
    x = Emu(900000) + i * Emu(5600000)
    y_base = Emu(3100000)

    # Cloud header bar
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y_base, Emu(5200000), Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_block(slide, x + Emu(200000), y_base + Emu(200000), Emu(4800000), Emu(500000),
                  cloud_name, FONT_BODY, Pt(16), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Fit items
    add_text_block(slide, x + Emu(100000), y_base + Emu(1000000), Emu(5000000), Emu(400000),
                  fit_label, FONT_CARD, Pt(12), color, bold=True)
    items_text = "\n".join([f"•  {item}" for item in fit_items])
    add_text_block(slide, x + Emu(100000), y_base + Emu(1400000), Emu(5000000), Emu(2200000),
                  items_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)

    # Use when
    add_text_block(slide, x + Emu(100000), y_base + Emu(3800000), Emu(5000000), Emu(400000),
                  use_label, FONT_CARD, Pt(12), color, bold=True)
    add_text_block(slide, x + Emu(100000), y_base + Emu(4200000), Emu(5000000), Emu(1200000),
                  use_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)

# Recommendation bar
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Emu(900000), Emu(8800000), Emu(16400000), Emu(900000))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_DK1
shape.line.fill.background()
add_text_block(slide, Emu(1200000), Emu(8900000), Emu(15800000), Emu(700000),
              "RECOMMENDATION: Start with Manufacturing Cloud as the commercial backbone (contracts + forecasting + rebates), "
              "layer Energy & Utilities patterns for asset-intensive operations, and adopt Consumer Goods visit execution for retail network management. "
              "Data Cloud unifies all three into one customer truth.",
              FONT_BODY, Pt(14), COLOR_WHITE, bold=False)


# ============================================================
# SLIDE 4: DATA CLOUD DEEP DIVE
# ============================================================
slide = add_content_slide(prs,
    "Data Cloud: The Unification Engine",
    "How Data Cloud creates a single customer, product, and operational truth across fragmented downstream systems"
)

# Data model entities
add_text_block(slide, MARGIN_LEFT, Emu(2300000), Emu(7500000), Emu(500000),
              "UNIFIED DATA MODEL FOR DOWNSTREAM:", FONT_CARD, Pt(14), COLOR_HEADER, bold=True)

entities = [
    "Account Profile — distributor, dealer, fleet, agribusiness, mining, industrial, TRR",
    "Site/Location — gas stations, terminals, bases, PAVs, on-site tanks, warehouses",
    "Asset — tanks (capacity, sensor), delivery equipment, fleet vehicles, infrastructure",
    "Product — diesel, gasoline, ethanol, biodiesel, lubricants, QAV, GLP, additives",
    "Contract — volume agreements, pricing terms, SLAs, credit limits, exclusivity",
    "Interaction — orders, deliveries, visits, cases, loyalty, payments, complaints",
    "Market Signal — Petrobras prices, FX, ANP data, competitor moves, demand indicators",
]
add_multi_text(slide, MARGIN_LEFT, Emu(2900000), Emu(8000000), Emu(5000000),
              [f"•  {e}" for e in entities], FONT_CARD, Pt(14), COLOR_CARD_BODY, spacing=Pt(6))

# Right side - capabilities
add_text_block(slide, Emu(9500000), Emu(2300000), Emu(7500000), Emu(500000),
              "KEY DATA CLOUD CAPABILITIES:", FONT_CARD, Pt(14), COLOR_HEADER, bold=True)

capabilities = [
    ("Identity Resolution", "Match dealer across SAP (vendor), CRM (account), loyalty (member), and fiscal (CNPJ) systems"),
    ("Calculated Insights", "Computed metrics: realized margin/account, delivery SLA, churn risk score, lifetime value"),
    ("Segmentation", "Dynamic segments: at-risk dealers, high-growth B2B, underperforming sites, loyalty-lapsed"),
    ("Streaming Ingestion", "Real-time signals: tank telemetry, price changes, FX moves, delivery events"),
    ("Activation", "Route insights to: Agentforce (recommendations), Tableau (dashboards), Sales Cloud (next best action)"),
]

y = Emu(2900000)
for i, (cap_name, cap_desc) in enumerate(capabilities):
    row_y = y + i * Emu(1200000)
    add_text_block(slide, Emu(9500000), row_y, Emu(7000000), Emu(400000),
                  cap_name, FONT_CARD, Pt(14), COLOR_DK2, bold=True)
    add_text_block(slide, Emu(9500000), row_y + Emu(400000), Emu(7000000), Emu(600000),
                  cap_desc, FONT_CARD, Pt(13), COLOR_CARD_BODY)

# Bottom insight
add_text_block(slide, MARGIN_LEFT, Emu(9200000), Emu(16000000), Emu(500000),
              "Architecture principle: Data Cloud does not replace SAP as system of record. It creates a decision-grade unified view that feeds analytics, agents, and engagement workflows.",
              FONT_BODY, Pt(14), COLOR_DK1, bold=True)
add_footer(slide)


# ============================================================
# SLIDE 5: MULESOFT INTEGRATION ARCHITECTURE
# ============================================================
slide = add_content_slide(prs,
    "MuleSoft: API-Led Connectivity for Downstream",
    "Three-tier API architecture connects 8+ source system categories to the Salesforce intelligence layer"
)

# Three API layers
api_layers = [
    ("SYSTEM APIs", COLOR_GREEN, "Connect to source systems",
     ["SAP IS-Oil API — master data, costs, inventory, credit, orders",
      "TMS API — routes, schedules, freight, delivery status, carrier events",
      "Fiscal API — NF-e emission, SEFAZ validation, tax calculations",
      "IoT/Telemetry API — tank levels, GPS, flow meters, quality sensors",
      "Market Data API — Petrobras prices, FX (B3), ANP indicators",
      "Loyalty/Cards API — transactions, points, fleet card authorization"]),
    ("PROCESS APIs", COLOR_DK2, "Compose business logic",
     ["Pricing Engine API — landed cost calculation, margin rules, corridor validation",
      "Replenishment API — demand signal + inventory → delivery recommendation",
      "Credit Decision API — exposure, score, limit, approval workflow trigger",
      "Dealer Health API — volume, loyalty, visits, cases, payment → health score",
      "Account Value API — margin, cost-to-serve, lifetime value, churn probability",
      "Fulfillment API — order → allocation → route → schedule → delivery event"]),
    ("EXPERIENCE APIs", COLOR_PURPLE, "Serve engagement channels",
     ["Dealer Portal API — orders, statements, stock visibility, promotions",
      "Fleet Portal API — consumption, controls, invoices, reporting",
      "Mobile Field API — visit plans, checklists, delivery proof, offline sync",
      "Agentforce API — context retrieval, action execution, feedback loop",
      "Tableau Embed API — dashboard context, drill-down, real-time refresh",
      "Notification API — alerts, approvals, escalations across channels"]),
]

for i, (layer_name, color, description, apis) in enumerate(api_layers):
    x = Emu(900000) + i * Emu(5600000)

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(2400000), Emu(5200000), Emu(900000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_block(slide, x + Emu(150000), Emu(2500000), Emu(4900000), Emu(400000),
                  layer_name, FONT_BODY, Pt(16), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_block(slide, x + Emu(150000), Emu(2900000), Emu(4900000), Emu(300000),
                  description, FONT_CONTENT, Pt(12), RGBColor(0xCC, 0xE0, 0xFF), alignment=PP_ALIGN.CENTER)

    # APIs list
    api_text = "\n".join([f"•  {api}" for api in apis])
    add_text_block(slide, x + Emu(100000), Emu(3500000), Emu(5100000), Emu(5000000),
                  api_text, FONT_CARD, Pt(12), COLOR_CARD_BODY)

# Governance bar
add_text_block(slide, MARGIN_LEFT, Emu(9000000), Emu(16400000), Emu(500000),
              "GOVERNANCE: API catalog • Version management • Rate limiting • Security policies • Monitoring • Reuse metrics",
              FONT_CARD, Pt(12), COLOR_MUTED)
add_text_block(slide, MARGIN_LEFT, Emu(9400000), Emu(16400000), Emu(400000),
              "Architecture principle: APIs are reusable assets. One System API serves multiple Process APIs. Process APIs serve any Experience API. No point-to-point integrations.",
              FONT_BODY, Pt(13), COLOR_DK1, bold=True)


# ============================================================
# SLIDE 6: AGENTFORCE ARCHITECTURE
# ============================================================
slide = add_content_slide(prs,
    "Agentforce: Trusted AI Agent Architecture",
    "How autonomous agents operate within governance boundaries across the five downstream engines"
)

# Agent anatomy
add_text_block(slide, MARGIN_LEFT, Emu(2200000), Emu(16000000), Emu(500000),
              "AGENT ANATOMY — How every downstream agent is structured:", FONT_CARD, Pt(14), COLOR_HEADER, bold=True)

# Flow: Trigger → Context → Reasoning → Action → Governance
flow_steps = [
    ("TRIGGER", COLOR_ORANGE, "Signal detected:\nprice move, volume\ndrop, SLA breach,\ncontract event"),
    ("CONTEXT", COLOR_DK2, "Data Cloud retrieves:\naccount, margin,\ncontract, service\nhistory, segment"),
    ("REASONING", COLOR_PURPLE, "Agent evaluates:\npolicy rules,\nthresholds, patterns,\nbest practices"),
    ("ACTION", COLOR_GREEN, "Recommendation:\nnext-best action,\nowner, urgency,\nrationale"),
    ("GOVERNANCE", COLOR_DK1, "Audit trail:\nsignal source, logic,\napproval path,\nKPI impact"),
]

for i, (step_name, color, description) in enumerate(flow_steps):
    x = Emu(900000) + i * Emu(3400000)
    y = Emu(2900000)

    # Step box
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(3100000), Emu(2200000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    add_text_block(slide, x + Emu(150000), y + Emu(150000), Emu(2800000), Emu(400000),
                  step_name, FONT_BODY, Pt(14), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    add_text_block(slide, x + Emu(150000), y + Emu(600000), Emu(2800000), Emu(1400000),
                  description, FONT_CARD, Pt(13), COLOR_WHITE)

    # Arrow between steps
    if i < len(flow_steps) - 1:
        add_text_block(slide, x + Emu(3100000), y + Emu(800000), Emu(300000), Emu(400000),
                      "→", FONT_BODY, Pt(22), COLOR_MUTED, bold=True)

# Agents deployed
add_text_block(slide, MARGIN_LEFT, Emu(5500000), Emu(16000000), Emu(500000),
              "AGENTS DEPLOYED ACROSS FIVE ENGINES:", FONT_CARD, Pt(14), COLOR_HEADER, bold=True)

agents_grid = [
    ("Pricing", ["Cost Monitor Agent", "Deal Margin Guardian", "Competitive Response Agent"]),
    ("Logistics", ["Predictive Replenishment Agent", "Exception Recovery Agent", "SLA Monitor Agent"]),
    ("Retail Network", ["Churn Prevention Agent", "Visit Priority Agent", "Compliance Alert Agent"]),
    ("B2B Growth", ["Opportunity Spotter Agent", "Deal Desk Agent", "Renewal Governance Agent"]),
    ("Customer Ops", ["Credit Risk Agent", "Loyalty Engagement Agent", "Fleet Expansion Agent"]),
]

for i, (engine, agent_list) in enumerate(agents_grid):
    x = Emu(900000) + i * Emu(3400000)
    y = Emu(6100000)

    add_text_block(slide, x, y, Emu(3200000), Emu(400000),
                  engine, FONT_CARD, Pt(13), COLOR_DK2, bold=True)
    agents_text = "\n".join([f"•  {a}" for a in agent_list])
    add_text_block(slide, x, y + Emu(450000), Emu(3200000), Emu(1500000),
                  agents_text, FONT_CARD, Pt(12), COLOR_CARD_BODY)

# Governance principle
add_text_block(slide, MARGIN_LEFT, Emu(9200000), Emu(16400000), Emu(500000),
              "TRUST PRINCIPLE: Agents recommend and document — they do not bypass human authority. Approval thresholds, margin floors, credit limits and compliance rules remain governed by policy.",
              FONT_BODY, Pt(13), COLOR_DK1, bold=True)


# ============================================================
# SLIDE 7: TABLEAU ANALYTICS ARCHITECTURE
# ============================================================
slide = add_content_slide(prs,
    "Tableau: Decision-Grade Analytics at Every Level",
    "From board-level KPIs to frontline operational views — one analytics platform serving all decision-makers"
)

# Analytics tiers
tiers = [
    ("EXECUTIVE / BOARD", COLOR_DK1, Emu(2400000),
     ["P&L by engine: realized margin, volume, cost-to-serve",
      "Portfolio health: network resilience, B2B pipeline, churn risk trajectory",
      "Transformation KPIs: adoption, agent performance, time-to-decision improvement",
      "Investment ROI: wave-by-wave value realization vs. baseline"]),
    ("MIDDLE MANAGEMENT / OPERATIONS", COLOR_DK2, Emu(4400000),
     ["Pricing: corridor compliance, exception rate, discount leakage by region/product",
      "Logistics: SLA adherence, route efficiency, stockout forecast, cost-per-delivery",
      "Commercial: pipeline velocity, visit execution, dealer health distribution, win rate",
      "Customer Ops: credit portfolio exposure, loyalty ROI, fleet utilization, case patterns"]),
    ("FRONTLINE / FIELD", COLOR_HEADER, Emu(6400000),
     ["Seller: my accounts today — who needs attention, what's the recommended action",
      "Field rep: my visit plan, station status, delivery schedule, open cases",
      "Pricing analyst: margin monitor, alerts queue, pending approvals, competitive signals",
      "Service agent: customer context, case priority, SLA countdown, resolution guidance"]),
]

for tier_name, color, y_pos, views in tiers:
    # Tier label bar
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   MARGIN_LEFT, y_pos, Emu(3000000), Emu(1600000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_block(slide, MARGIN_LEFT + Emu(150000), y_pos + Emu(400000),
                  Emu(2700000), Emu(800000),
                  tier_name, FONT_BODY, Pt(13), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)

    # Views
    views_text = "\n".join([f"•  {v}" for v in views])
    add_text_block(slide, Emu(4200000), y_pos + Emu(100000), Emu(13000000), Emu(1500000),
                  views_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)

# Connection to agents
add_text_block(slide, MARGIN_LEFT, Emu(8600000), Emu(16400000), Emu(400000),
              "TABLEAU + AGENTFORCE: Tableau Pulse surfaces anomalies and KPI shifts. Agentforce interprets the context and recommends action. Together they create 'insight-to-action' loops.",
              FONT_BODY, Pt(14), COLOR_DK1, bold=True)
add_text_block(slide, MARGIN_LEFT, Emu(9200000), Emu(16400000), Emu(400000),
              "DATA SOURCE: Tableau connects to Data Cloud (unified truth), SAP (transactional detail), and TMS (operational metrics) through governed data models.",
              FONT_CARD, Pt(12), COLOR_MUTED)


# ============================================================
# SLIDE 8: INTEGRATION LANDSCAPE (Systems Map)
# ============================================================
slide = add_content_slide(prs,
    "Source System Landscape: What MuleSoft Connects",
    "Typical system map for a Brazilian fuel distributor — 8 integration domains, 20+ system touchpoints"
)

systems = [
    ("ERP / Finance", COLOR_DK1,
     ["SAP S/4HANA or ECC", "SAP IS-Oil&Gas", "SAP MM, FI, CO-PA, SD", "SAP BW/4HANA (reporting)"],
     "Master data • Costs • Orders\nInventory • Credit • GL posting"),
    ("Transport & Logistics", COLOR_DK2,
     ["SAP TM or custom TMS", "Logix / TOTVS Logística", "Fleet GPS systems", "Port/terminal systems"],
     "Routes • Freight • Schedules\nCarrier allocation • POD"),
    ("Fiscal & Compliance", COLOR_GREEN,
     ["SEFAZ (NF-e, CT-e)", "SPED/EFD systems", "ICMS calculation engine", "CBIOs/RenovaBio platform"],
     "Tax documents • State rules\nCompliance • Environment"),
    ("Pricing & Market", COLOR_ORANGE,
     ["Petrobras price portal", "ANP open data APIs", "B3 / Bloomberg (FX)", "Internal pricing engine"],
     "Gate prices • Benchmarks\nFX rates • Competitor data"),
    ("IoT & Telemetry", COLOR_PURPLE,
     ["Tank level sensors (Veeder-Root)", "SCADA systems", "Fleet telematics", "Flow meters / quality"],
     "Tank levels • Vehicle GPS\nConsumption • Quality"),
    ("Loyalty & Payments", RGBColor(0xC6, 0x28, 0x28),
     ["Card processor / acquirer", "Legacy loyalty platform", "Fleet card system", "Digital wallet (app)"],
     "Transactions • Points\nAuthorization • Redemption"),
    ("HR & Workforce", COLOR_MUTED,
     ["SAP HCM / SuccessFactors", "Shift management", "Training / certification", "Safety systems"],
     "Field teams • Drivers\nCompliance • Safety"),
    ("External / Partners", RGBColor(0x00, 0x7B, 0x5F),
     ["Carrier EDI / portals", "Dealer management legacy", "Serasa/Bureau credit", "Government (ANP, ANTT)"],
     "Partner capacity • Credit\nRegulatory • Certification"),
]

for i, (category, color, platforms, data_desc) in enumerate(systems):
    col = i % 4
    row = i // 4
    x = Emu(900000) + col * Emu(4200000)
    y = Emu(2400000) + row * Emu(3600000)

    # Card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(3900000), Emu(3200000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    # Color accent bar at top
    shape2 = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, Emu(3900000), Emu(80000))
    shape2.fill.solid()
    shape2.fill.fore_color.rgb = color
    shape2.line.fill.background()

    # Category
    add_text_block(slide, x + Emu(150000), y + Emu(200000), Emu(3600000), Emu(400000),
                  category, FONT_CARD, Pt(14), color, bold=True)
    # Platforms
    plat_text = "\n".join([f"  {p}" for p in platforms])
    add_text_block(slide, x + Emu(150000), y + Emu(650000), Emu(3600000), Emu(1400000),
                  plat_text, FONT_CARD, Pt(12), COLOR_CARD_BODY)
    # Data description
    add_text_block(slide, x + Emu(150000), y + Emu(2200000), Emu(3600000), Emu(800000),
                  f"→ {data_desc}", FONT_CARD, Pt(11), COLOR_MUTED)

# Bottom
add_text_block(slide, MARGIN_LEFT, Emu(9500000), Emu(16400000), Emu(400000),
              "Architecture principle: MuleSoft creates reusable APIs from each domain. Each API is independently versioned, monitored, and governed. No direct system-to-system coupling.",
              FONT_BODY, Pt(12), COLOR_DK1, bold=True)


# ============================================================
# SLIDE 9: DATA FLOW — END TO END SCENARIO
# ============================================================
slide = add_content_slide(prs,
    "End-to-End Data Flow: From Signal to Action",
    "Tracing a real scenario through the full architecture stack"
)

# Scenario description
add_text_block(slide, MARGIN_LEFT, Emu(2200000), Emu(16000000), Emu(500000),
              "SCENARIO: Agribusiness account in Mato Grosso shows volume decline + delivery complaint + competitor activity detected",
              FONT_BODY, Pt(16), COLOR_DK2, bold=True)

# Flow steps
steps = [
    ("1. SIGNAL DETECTION", "IoT + SAP + CRM",
     "Tank telemetry shows 30% less consumption vs. last month\nSAP shows reduced order frequency\nField rep logged competitor visit at site"),
    ("2. DATA UNIFICATION", "MuleSoft → Data Cloud",
     "MuleSoft APIs pull SAP order history + telemetry + CRM activity\nData Cloud resolves identity, calculates churn score = 78%\nSegment: high-value agribusiness, harvest season approaching"),
    ("3. INSIGHT GENERATION", "Tableau + Einstein",
     "Tableau Pulse flags account in 'critical risk' segment\nEinstein model: 78% probability of switching within 60 days\nMargin at risk: R$1.2M annual volume"),
    ("4. AGENT RECOMMENDATION", "Agentforce",
     "Churn Prevention Agent activates\nContext: harvest season = high-volume period approaching\nRecommendation: executive visit + competitive pricing review + delivery SLA upgrade"),
    ("5. GOVERNED EXECUTION", "Sales + Service + Field",
     "Task created for regional manager (urgent visit within 7 days)\nPricing team receives corridor review request with margin context\nField Service schedules proactive delivery for next harvest window"),
    ("6. MEASUREMENT", "Tableau + Data Cloud",
     "Track: account saved/lost, volume recovery, margin outcome\nFeedback loop: agent reasoning quality, policy effectiveness\nLearning: update churn model with outcome data"),
]

y = Emu(2900000)
for i, (step_name, system, description) in enumerate(steps):
    row_y = y + i * Emu(1100000)

    # Step number/name
    add_text_block(slide, MARGIN_LEFT, row_y, Emu(3200000), Emu(400000),
                  step_name, FONT_CARD, Pt(13), COLOR_DK2, bold=True)
    # System
    add_text_block(slide, Emu(4300000), row_y, Emu(2500000), Emu(400000),
                  system, FONT_CARD, Pt(12), COLOR_MUTED, bold=True)
    # Description
    add_text_block(slide, Emu(7000000), row_y, Emu(10000000), Emu(900000),
                  description, FONT_CARD, Pt(12), COLOR_CARD_BODY)

add_footer(slide)


# ============================================================
# SLIDE 10: SECURITY & GOVERNANCE
# ============================================================
slide = add_content_slide(prs,
    "Security, Governance & Compliance",
    "Enterprise-grade controls for regulated downstream operations"
)

# Four domains
domains = [
    ("Data Governance", COLOR_DK2,
     ["Data classification: public, internal, confidential, restricted",
      "Field-level security: margin data visible only to pricing + leadership",
      "Record sharing: territory-based for dealers, role-based for B2B",
      "Retention policies: aligned with ANP and fiscal requirements (5+ years)",
      "LGPD compliance: consent management for consumer loyalty data"]),
    ("Integration Security", COLOR_GREEN,
     ["MuleSoft API gateway: OAuth 2.0, mTLS, rate limiting",
      "API-level RBAC: each consumer gets scoped credentials",
      "Event encryption: at-rest and in-transit for all data flows",
      "Audit logging: every API call tracked for compliance",
      "Network isolation: VPC peering for SAP, private endpoints"]),
    ("AI Governance", COLOR_PURPLE,
     ["Agentforce trust layer: no customer data used for model training",
      "Agent audit trail: every recommendation logged with reasoning",
      "Human-in-the-loop: configurable approval thresholds by risk level",
      "Explainability: agents document why they recommended each action",
      "Bias monitoring: periodic review of agent outcomes by segment"]),
    ("Operational Compliance", COLOR_ORANGE,
     ["ANP regulatory reporting: automated data extraction + submission",
      "CBIOs / RenovaBio: credit tracking and compliance workflows",
      "SEFAZ integration: real-time fiscal document validation",
      "Audit readiness: full trail from signal to decision to action",
      "SOX controls: segregation of duties in approval workflows"]),
]

for i, (domain_name, color, items) in enumerate(domains):
    col = i % 2
    row = i // 2
    x = Emu(900000) + col * Emu(8400000)
    y = Emu(2400000) + row * Emu(3600000)

    # Header
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(7800000), Emu(700000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_block(slide, x + Emu(200000), y + Emu(150000), Emu(7400000), Emu(400000),
                  domain_name, FONT_BODY, Pt(15), COLOR_WHITE, bold=True)

    # Items
    items_text = "\n".join([f"•  {item}" for item in items])
    add_text_block(slide, x + Emu(200000), y + Emu(850000), Emu(7400000), Emu(2500000),
                  items_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)

add_text_block(slide, MARGIN_LEFT, Emu(9600000), Emu(16400000), Emu(300000),
              "Salesforce Trust: SOC 2 Type II • ISO 27001 • LGPD ready • Data residency options (Brazil) • Encryption at rest & transit • Shield Platform Encryption available",
              FONT_CARD, Pt(11), COLOR_MUTED)


# ============================================================
# SLIDE 11: DEPLOYMENT & ENVIRONMENT STRATEGY
# ============================================================
slide = add_content_slide(prs,
    "Deployment Strategy: Environments & Release Governance",
    "How to scale from MVP to enterprise with controlled releases and quality gates"
)

# Environment stack
envs = [
    ("PRODUCTION", COLOR_DK1, "Live operations — governed releases only"),
    ("STAGING / UAT", COLOR_DK2, "Business validation — user acceptance before promotion"),
    ("QA / INTEGRATION", COLOR_HEADER, "Integration testing — full system connectivity validation"),
    ("DEVELOPMENT", COLOR_GREEN, "Build & unit test — developer sandboxes + scratch orgs"),
]

for i, (env_name, color, desc) in enumerate(envs):
    y = Emu(2400000) + i * Emu(1000000)
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                                   MARGIN_LEFT, y, Emu(7500000), Emu(800000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    add_text_block(slide, MARGIN_LEFT + Emu(200000), y + Emu(100000), Emu(3000000), Emu(400000),
                  env_name, FONT_BODY, Pt(14), COLOR_WHITE, bold=True)
    add_text_block(slide, MARGIN_LEFT + Emu(200000), y + Emu(420000), Emu(7000000), Emu(300000),
                  desc, FONT_CONTENT, Pt(12), RGBColor(0xCC, 0xE0, 0xFF))

# Right side - release governance
add_text_block(slide, Emu(9500000), Emu(2400000), Emu(7500000), Emu(500000),
              "RELEASE GOVERNANCE:", FONT_CARD, Pt(14), COLOR_HEADER, bold=True)

release_items = [
    "CI/CD pipeline: Salesforce DevOps Center or Copado",
    "Version control: Git-based metadata management",
    "Quality gates: code review, test coverage >80%, regression suite",
    "Release cadence: bi-weekly sprints, monthly production releases",
    "Rollback plan: feature flags + metadata snapshots",
    "Change Advisory Board: business + IT approval for major releases",
    "",
    "TEAM MODEL:",
    "Platform team: Salesforce architecture, DevOps, security",
    "Engine squads: dedicated team per business engine (Pricing, Logistics, etc.)",
    "Integration team: MuleSoft APIs, connectors, event architecture",
    "Analytics team: Tableau dashboards, Data Cloud models, Einstein",
    "AI/Agent team: Agentforce configuration, testing, governance",
]
add_multi_text(slide, Emu(9500000), Emu(3000000), Emu(7500000), Emu(6500000),
              release_items, FONT_CARD, Pt(13), COLOR_CARD_BODY, spacing=Pt(5))

add_footer(slide)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/iaraujo/Downloads/OilGas_Downstream_POV_Architecture_Deep.pptx"
prs.save(output_path)
print(f"✓ Saved: {output_path}")
print(f"  Total slides: {len(prs.slides)}")
