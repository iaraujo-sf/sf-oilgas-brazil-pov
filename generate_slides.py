"""
Generate Oil&Gas Downstream POV expanded sections as .pptx
Matching Salesforce corporate template visual style
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu, Cm
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import copy

# --- CONSTANTS matching original template ---
SLIDE_WIDTH = Emu(18288000)   # 20 inches
SLIDE_HEIGHT = Emu(10287000)  # 11.25 inches

# Colors from theme
COLOR_DK1 = RGBColor(0x00, 0x1E, 0x5B)      # Dark blue (titles)
COLOR_DK2 = RGBColor(0x02, 0x2A, 0xC0)      # Electric blue
COLOR_LT2 = RGBColor(0x00, 0xB3, 0xFF)      # Light blue accent
COLOR_ACCENT1 = RGBColor(0x90, 0xD0, 0xFE)  # Cloud blue
COLOR_ACCENT2 = RGBColor(0x04, 0xE1, 0xCB)  # Teal
COLOR_ACCENT3 = RGBColor(0xFC, 0xC0, 0x03)  # Yellow
COLOR_ACCENT6 = RGBColor(0xEA, 0xF5, 0xFE)  # Very light blue bg

# Content colors from analysis
COLOR_HEADER = RGBColor(0x12, 0x37, 0x7A)    # Card headers
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)      # Body text
COLOR_MUTED = RGBColor(0x85, 0x84, 0x81)     # Footer/muted
COLOR_CARD_BODY = RGBColor(0x38, 0x46, 0x53) # Card body (from slide 25)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_RED = RGBColor(0xCC, 0x00, 0x00)
COLOR_ORANGE = RGBColor(0xE6, 0x8A, 0x00)
COLOR_BG_LIGHT = RGBColor(0xCF, 0xE9, 0xFE) # Cloud Blue 90

# Fonts
FONT_TITLE = "Avant Garde Demi SFDC"
FONT_BODY = "Salesforce Sans"
FONT_CARD = "Inter"
FONT_CONTENT = "Montserrat"

# Margins
MARGIN_LEFT = Emu(900000)   # ~1 inch
MARGIN_TOP = Emu(600000)
CONTENT_TOP = Emu(2000000)  # Below title area


def create_presentation():
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    return prs


def add_section_divider(prs, title, subtitle):
    """Dark blue background section divider"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Dark blue background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DK1

    # Title
    txBox = slide.shapes.add_textbox(
        Emu(1800000), Emu(3200000), Emu(14400000), Emu(2400000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(44)
    run.font.color.rgb = COLOR_WHITE
    run.font.bold = False

    # Subtitle
    txBox2 = slide.shapes.add_textbox(
        Emu(1800000), Emu(5800000), Emu(14400000), Emu(1600000)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.alignment = PP_ALIGN.LEFT
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_BODY
    run2.font.size = Pt(22)
    run2.font.color.rgb = COLOR_ACCENT1
    run2.font.bold = False

    return slide


def add_content_slide(prs, title, subtitle=None):
    """Standard content slide with white background"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

    # Light background
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    # Title
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, MARGIN_TOP, Emu(16000000), Emu(1200000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(36)
    run.font.color.rgb = COLOR_DK1
    run.font.bold = False

    # Subtitle if provided
    if subtitle:
        txBox2 = slide.shapes.add_textbox(
            MARGIN_LEFT, Emu(1400000), Emu(16000000), Emu(800000)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_BODY
        run2.font.size = Pt(20)
        run2.font.color.rgb = COLOR_BODY

    return slide


def add_text_block(slide, left, top, width, height, text, font_name=FONT_BODY,
                   font_size=Pt(16), color=COLOR_BODY, bold=False, alignment=PP_ALIGN.LEFT):
    """Add a text block to a slide"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return tf


def add_multi_text(slide, left, top, width, height, lines, font_name=FONT_BODY,
                   font_size=Pt(16), color=COLOR_BODY):
    """Add multiple lines of text"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(line, tuple):
            text, is_bold, line_color, size = line
            run = p.add_run()
            run.text = text
            run.font.name = font_name
            run.font.size = size if size else font_size
            run.font.color.rgb = line_color if line_color else color
            run.font.bold = is_bold
        else:
            run = p.add_run()
            run.text = line
            run.font.name = font_name
            run.font.size = font_size
            run.font.color.rgb = color
            run.font.bold = False

        p.space_after = Pt(6)

    return tf


def add_card(slide, left, top, width, height, header, body, header_color=COLOR_HEADER):
    """Add a card with header and body text"""
    # Card background
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    # Header
    txBox = slide.shapes.add_textbox(
        left + Emu(200000), top + Emu(200000),
        width - Emu(400000), Emu(600000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = header
    run.font.name = FONT_CARD
    run.font.size = Pt(16)
    run.font.color.rgb = header_color
    run.font.bold = True

    # Body
    txBox2 = slide.shapes.add_textbox(
        left + Emu(200000), top + Emu(800000),
        width - Emu(400000), height - Emu(1000000)
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = body
    run2.font.name = FONT_CARD
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLOR_CARD_BODY

    return shape


def add_table_slide(prs, title, subtitle, headers, rows):
    """Add a slide with a table"""
    slide = add_content_slide(prs, title, subtitle)

    cols = len(headers)
    table_rows = len(rows) + 1  # +1 for header

    # Table position
    left = Emu(900000)
    top = Emu(2600000)
    width = Emu(16400000)
    height = Emu(6800000)
    row_height = height // table_rows

    table_shape = slide.shapes.add_table(table_rows, cols, left, top, width, height)
    table = table_shape.table

    # Set column widths proportionally
    col_width = width // cols
    for i in range(cols):
        table.columns[i].width = col_width

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = COLOR_DK1
        para = cell.text_frame.paragraphs[0]
        para.font.name = FONT_BODY
        para.font.size = Pt(14)
        para.font.color.rgb = COLOR_WHITE
        para.font.bold = True

    # Data rows
    for r_idx, row in enumerate(rows):
        for c_idx, cell_text in enumerate(row):
            cell = table.cell(r_idx + 1, c_idx)
            cell.text = cell_text
            # Alternate row colors
            if r_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_ACCENT6
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = COLOR_WHITE
            para = cell.text_frame.paragraphs[0]
            para.font.name = FONT_CARD
            para.font.size = Pt(13)
            para.font.color.rgb = COLOR_BODY

    return slide


def add_gap_slide(prs, title, subtitle, layers):
    """Add an operating gap slide (matching slide 24 style)"""
    slide = add_content_slide(prs, title, subtitle)

    y_start = Emu(2600000)
    row_height = Emu(1300000)

    for i, (layer_name, current_reality, risk_text, risk_color) in enumerate(layers):
        y = y_start + (i * row_height)

        # Layer number + name
        add_text_block(slide, Emu(900000), y, Emu(5000000), Emu(500000),
                      f"{i+1} · {layer_name}", FONT_BODY, Pt(22), COLOR_DK1, bold=True)

        # Current reality
        add_text_block(slide, Emu(900000), y + Emu(500000), Emu(8000000), Emu(400000),
                      current_reality, FONT_BODY, Pt(17), COLOR_BODY)

        # Risk indicator
        color = COLOR_RED if risk_color == "red" else COLOR_ORANGE
        add_text_block(slide, Emu(12000000), y + Emu(200000), Emu(5500000), Emu(600000),
                      risk_text, FONT_BODY, Pt(14), color, bold=True)

    return slide


def add_footer(slide, text="Salesforce POV • Oil&Gas Downstream"):
    """Add footer text"""
    txBox = slide.shapes.add_textbox(
        MARGIN_LEFT, Emu(9600000), Emu(8000000), Emu(400000)
    )
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CARD
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_MUTED


def add_quote_slide(prs, quote, attribution=None):
    """Add a quote/insight slide"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])

    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_BG_LIGHT

    txBox = slide.shapes.add_textbox(
        Emu(2400000), Emu(3000000), Emu(13000000), Emu(4000000)
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f'"{quote}"'
    run.font.name = FONT_BODY
    run.font.size = Pt(28)
    run.font.color.rgb = COLOR_DK1
    run.font.bold = False

    if attribution:
        p2 = tf.add_paragraph()
        p2.space_before = Pt(24)
        run2 = p2.add_run()
        run2.text = f"— {attribution}"
        run2.font.name = FONT_BODY
        run2.font.size = Pt(18)
        run2.font.color.rgb = COLOR_MUTED

    return slide


def add_mapping_slide(prs, title, subtitle, layers):
    """Add a Salesforce mapping slide (matching slide 25 style)"""
    slide = add_content_slide(prs, title, subtitle)

    y_start = Emu(2600000)
    col_width = Emu(3200000)
    num_layers = len(layers)

    for i, (product, capabilities) in enumerate(layers):
        x = Emu(900000) + (i * (col_width + Emu(200000)))

        # Product header card
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y_start, col_width, Emu(1000000)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DK2
        shape.line.fill.background()

        # Product name
        txBox = slide.shapes.add_textbox(
            x + Emu(100000), y_start + Emu(200000),
            col_width - Emu(200000), Emu(600000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = product
        run.font.name = FONT_BODY
        run.font.size = Pt(15)
        run.font.color.rgb = COLOR_WHITE
        run.font.bold = True

        # Capabilities
        cap_y = y_start + Emu(1200000)
        for j, cap in enumerate(capabilities):
            txBox2 = slide.shapes.add_textbox(
                x + Emu(100000), cap_y + (j * Emu(700000)),
                col_width - Emu(200000), Emu(650000)
            )
            tf2 = txBox2.text_frame
            tf2.word_wrap = True
            p2 = tf2.paragraphs[0]
            run2 = p2.add_run()
            run2.text = cap
            run2.font.name = FONT_CONTENT
            run2.font.size = Pt(14)
            run2.font.color.rgb = COLOR_CARD_BODY

    return slide


def add_agentforce_slide(prs, title, agents):
    """Add Agentforce slide with agent cards"""
    slide = add_content_slide(prs, title, "Autonomous Intelligence with Agentforce")

    num_agents = len(agents)
    card_width = Emu(16000000) // min(num_agents, 3)
    card_gap = Emu(200000)

    for i, (agent_name, bullets) in enumerate(agents):
        row = i // 3
        col = i % 3
        x = Emu(900000) + col * (card_width + card_gap)
        y = Emu(2600000) + row * Emu(3800000)

        # Agent card background
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, x, y,
            card_width - card_gap, Emu(3500000)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_ACCENT6
        shape.line.fill.background()

        # Agent name
        txBox = slide.shapes.add_textbox(
            x + Emu(200000), y + Emu(200000),
            card_width - Emu(600000), Emu(600000)
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        run = p.add_run()
        run.text = agent_name
        run.font.name = FONT_CARD
        run.font.size = Pt(16)
        run.font.color.rgb = COLOR_DK2
        run.font.bold = True

        # Bullets
        txBox2 = slide.shapes.add_textbox(
            x + Emu(200000), y + Emu(900000),
            card_width - Emu(600000), Emu(2400000)
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        for j, bullet in enumerate(bullets):
            if j == 0:
                p2 = tf2.paragraphs[0]
            else:
                p2 = tf2.add_paragraph()
            p2.space_after = Pt(6)
            run2 = p2.add_run()
            run2.text = f"• {bullet}"
            run2.font.name = FONT_CARD
            run2.font.size = Pt(13)
            run2.font.color.rgb = COLOR_CARD_BODY

    add_footer(slide)
    return slide


# ============================================================
# BUILD THE PRESENTATION
# ============================================================

prs = create_presentation()

# ============================================================
# SECTION 1: PRICING EXCELLENCE
# ============================================================

add_section_divider(prs,
    "Margin Through Pricing Excellence",
    "How leading distributors protect profitability while remaining competitive across channels and regions."
)

# Slide: Why Pricing Is Different
slide = add_content_slide(prs,
    "Pricing in Brazil: A Multi-Variable Equation Unlike Any Other Market",
    "Five structural factors make fuel pricing in Brazil uniquely complex"
)
add_table_slide(prs,
    "Pricing in Brazil: A Multi-Variable Equation",
    "Five structural factors make fuel pricing in Brazil uniquely complex",
    ["Factor", "Impact", "Pricing Implication"],
    [
        ["PPI Abandonment (2022)", "Petrobras no longer follows International Parity Price — opacity in cost basis", "Distributors must build own cost models"],
        ["FX Exposure", "BRL/USD volatility directly impacts import costs", "Timing risk between purchase and resale"],
        ["Tax Labyrinth", "ICMS monofásico + PIS/COFINS + CIDE — each state/product combination differs", "Tax intelligence = margin driver"],
        ["Biofuel Mandates", "Ethanol (27%) + Biodiesel (15%) blend ratios shift effective cost", "Pricing must reflect blend economics"],
        ["Competitive Pressure", "White-flag stations undercut branded networks on price", "Constant volume vs. margin trade-off"],
    ]
)

# Slide: Pricing Operating Model
slide = add_content_slide(prs,
    "The Pricing Decision Architecture",
    "Five interconnected layers that determine margin capture — or margin leakage."
)
layers_data = [
    ("Layer 1: Cost Basis Construction", [
        "Petrobras gate price (or import landed cost)",
        "Biofuel blending cost (ethanol + biodiesel spot)",
        "Freight to terminal + storage costs",
        "Tax burden per state/product/regime"
    ]),
    ("Layer 2: Market Intelligence", [
        "Competitor pricing monitoring (ANP data)",
        "Regional demand elasticity",
        "Channel-specific willingness to pay"
    ]),
    ("Layer 3: Margin Rules Engine", [
        "Floor margin by product/region/channel",
        "Target margin bands by customer tier",
        "Escalation triggers below threshold"
    ]),
    ("Layer 4: Dynamic Pricing Execution", [
        "B2B contracts (indexed to Petrobras + FX)",
        "Spot pricing for TRR/uncontracted volume",
        "Promotional pricing for retail network"
    ]),
    ("Layer 5: Performance Monitoring", [
        "Realized margin vs. target (daily)",
        "Price-volume trade-off analysis",
        "Competitor response tracking"
    ])
]

y = Emu(2600000)
for i, (layer_name, items) in enumerate(layers_data):
    x = Emu(900000) + (i * Emu(3400000))

    # Layer header
    add_text_block(slide, x, y, Emu(3200000), Emu(600000),
                  layer_name, FONT_CARD, Pt(14), COLOR_HEADER, bold=True)

    # Items
    item_text = "\n".join([f"• {item}" for item in items])
    add_text_block(slide, x, y + Emu(650000), Emu(3200000), Emu(2800000),
                  item_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)

add_footer(slide)

# Slide: Pricing Gap
add_gap_slide(prs,
    "The Gap: World-Class Complexity, Spreadsheet-Grade Tools",
    "Pricing decisions are made with yesterday's data in a market that changes hourly.",
    [
        ("Cost Basis Construction", "SAP extract + manual FX lookup + tax team email", "🟠 MARGIN EROSION — 24-48h delay in cost visibility", "orange"),
        ("Market Intelligence", "Manual pump surveys + weekly ANP reports + WhatsApp groups", "🟠 BLIND SPOTS — competitors react faster", "orange"),
        ("Margin Rules Engine", "Pricing committee meets weekly + Excel guardrails", "🔴 MARGIN LEAKAGE — deals close below floor", "red"),
        ("Dynamic Pricing Execution", "Sales rep discretion + email approval chains", "🔴 SPEED LOSS — B2B deals take days not hours", "red"),
        ("Performance Monitoring", "Monthly P&L analysis + ad-hoc reports", "🔴 LATE DETECTION — problems found weeks later", "red"),
    ]
)

# Slide: Pricing Salesforce Mapping
add_table_slide(prs,
    "The Five Pricing Layers — Salesforce Mapping",
    "Each layer powered by a specific Salesforce product",
    ["Layer", "Salesforce Product", "Key Capability"],
    [
        ["Cost Basis Construction", "Data Cloud + MuleSoft", "Real-time SAP cost ingestion + FX feeds + tax engine. Single cost-per-liter view updated continuously"],
        ["Market Intelligence", "Tableau + Data Cloud", "ANP data integration, competitor dashboards, regional elasticity models. AI-detected pricing anomalies"],
        ["Margin Rules Engine", "Revenue Cloud (CPQ) + Flow", "Automated floor/ceiling guardrails per product/region/channel. No deal closes below threshold without VP override"],
        ["Dynamic Pricing Execution", "Revenue Cloud + Agentforce", "AI agent recommends optimal price per B2B opportunity. Auto-generates indexed contract terms. Sub-hour response"],
        ["Performance Monitoring", "CRM Analytics + Tableau", "Daily realized margin dashboards. Automated alerts when margin drifts. Price-volume elasticity tracking"],
    ]
)

# Slide: Pricing Agentforce
add_agentforce_slide(prs,
    "Pricing — Agentforce in Action",
    [
        ("Cost Monitor Agent", [
            "Watches Petrobras gate price changes, FX movements, biofuel spot prices",
            "Calculates new landed cost within minutes of market move",
            "Triggers pricing review alert with recommended new price bands"
        ]),
        ("Deal Margin Guardian", [
            "Reviews every B2B deal before approval",
            "Validates margin against floor rules",
            "Auto-approves within guardrails, escalates exceptions with context"
        ]),
        ("Competitive Response Agent", [
            "Monitors regional price movements from ANP data and field intel",
            "Detects when competitor undercuts in specific regions",
            "Recommends targeted response based on account stickiness"
        ]),
    ]
)

# Quote
add_quote_slide(prs,
    "The goal is not to replace the pricing committee — it's to give them real-time ammunition instead of last week's spreadsheet.",
    "Pricing Excellence Vision"
)

# ============================================================
# SECTION 2: LOGISTICS EXCELLENCE
# ============================================================

add_section_divider(prs,
    "Scale Through Logistics Excellence",
    "How leading players turn network efficiency and service reliability into commercial advantage."
)

# Slide: Why Logistics Matters
add_table_slide(prs,
    "Logistics in Brazil: Where Continental Scale Meets Operational Reality",
    "Logistics cost = 8-12% of total fuel cost. Every 1% efficiency gain = R$500M+ industry savings annually.",
    ["Factor", "Brazil Context", "Implication"],
    [
        ["Continental Distances", "Average delivery route: 300-800km. SP to Manaus: 3,900km by road", "Freight cost dominates landed cost in remote regions"],
        ["Modal Dependence", "65% of fuel transported by road (tanker trucks) — highest cost modal", "Limited pipeline/rail constrains optimization"],
        ["Infrastructure Gaps", "Limited pipeline coverage outside Southeast. Port congestion", "Demand for primary storage & import terminal capacity"],
        ["Fleet Complexity", "Mix of owned, contracted, and third-party carriers", "Different SLAs, costs, and control levels"],
        ["Delivery Precision", "44,000+ stations expect delivery within tight windows", "Stockout = lost sales + customer churn to competitors"],
    ]
)

# Slide: Logistics Operating Model
slide = add_content_slide(prs,
    "The Logistics Decision Architecture",
    "From terminal to station — five layers that determine cost, reliability, and competitive advantage."
)
log_layers = [
    ("Layer 1: Network Design", [
        "Terminal location strategy",
        "Pipeline access contracts",
        "Import terminal capacity"
    ]),
    ("Layer 2: Primary Transport", [
        "Pipeline scheduling",
        "Cabotage (coastal shipping)",
        "Rail + long-haul road"
    ]),
    ("Layer 3: Secondary Transport", [
        "Route optimization per window",
        "Fleet mix allocation",
        "Multi-product compartment loads"
    ]),
    ("Layer 4: Last-Mile Execution", [
        "Delivery scheduling via tank levels",
        "Driver safety & compliance",
        "Proof of delivery & reconciliation"
    ]),
    ("Layer 5: Cost Management", [
        "Cost/liter by route/modal/region",
        "Fleet utilization rate",
        "SLA compliance & empty-mile reduction"
    ])
]
y = Emu(2600000)
for i, (layer_name, items) in enumerate(log_layers):
    x = Emu(900000) + (i * Emu(3400000))
    add_text_block(slide, x, y, Emu(3200000), Emu(600000),
                  layer_name, FONT_CARD, Pt(14), COLOR_HEADER, bold=True)
    item_text = "\n".join([f"• {item}" for item in items])
    add_text_block(slide, x, y + Emu(650000), Emu(3200000), Emu(2400000),
                  item_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)
add_footer(slide)

# Slide: Logistics Gap
add_gap_slide(prs,
    "The Gap: Millions of KM Managed on Calls and Spreadsheets",
    "Logistics decisions are made independently from commercial priorities.",
    [
        ("Network Design", "Annual planning cycle + static models", "🟠 SUBOPTIMAL CAPEX — infrastructure lags demand", "orange"),
        ("Primary Transport", "TMS scheduling disconnected from commercial demand", "🟠 EXCESS COST — over/under positioning at bases", "orange"),
        ("Secondary Transport", "Manual route planning + driver experience", "🔴 COST LEAKAGE — 15-25% empty miles average", "red"),
        ("Last-Mile Execution", "Phone calls to stations + driver judgment", "🔴 STOCKOUT RISK — reactive delivery triggers emergency runs", "red"),
        ("Cost Management", "Monthly reports from TMS + manual freight audit", "🟠 LATE VISIBILITY — cost overruns detected after the fact", "orange"),
    ]
)

# Slide: Logistics Salesforce Mapping
add_table_slide(prs,
    "The Five Logistics Layers — Salesforce Mapping",
    "Connecting physical assets to commercial decision-making",
    ["Layer", "Salesforce Product", "Key Capability"],
    [
        ["Network Design", "Tableau + Data Cloud", "Demand heat maps overlaid with infrastructure. AI-driven capacity planning. Investment case generation"],
        ["Primary Transport", "MuleSoft + Data Cloud", "TMS integration. Real-time pipeline/vessel tracking. Demand-triggered positioning recommendations"],
        ["Secondary Transport", "Field Service + Agentforce", "AI-optimized routing. Dynamic scheduling based on tank telemetry. Fleet allocation by priority"],
        ["Last-Mile Execution", "Field Service + Mobile", "Real-time delivery tracking. Automated proof of delivery. Driver tasks. Customer notifications"],
        ["Cost Management", "Tableau + CRM Analytics", "Cost-per-liter dashboards by route/region/customer. Freight audit automation. SLA compliance tracking"],
    ]
)

# Slide: Logistics Connected Intelligence
slide = add_content_slide(prs,
    "From Reactive Delivery to Predictive Logistics",
    "Scenario: Station Tank Telemetry → Autonomous Delivery Scheduling"
)
flow_steps = [
    ("IoT Sensor at Station", "Tank level drops below 40%"),
    ("Data Cloud", "Ingests telemetry + correlates with historical consumption"),
    ("Agentforce", "Predicts stockout in 36h. Checks delivery schedule."),
    ("Field Service", "Auto-schedules delivery. Optimizes route with nearby stations."),
    ("Experience Cloud", "Dealer notified: 'Delivery confirmed for tomorrow 8-10am'"),
    ("CRM Analytics", "Updates cost-per-delivery and SLA metrics in real time"),
]
y = Emu(2800000)
for i, (system, action) in enumerate(flow_steps):
    row_y = y + (i * Emu(1050000))
    # System badge
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(900000), row_y, Emu(3600000), Emu(800000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_DK2
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(Emu(1000000), row_y + Emu(200000), Emu(3400000), Emu(400000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = system
    run.font.name = FONT_BODY
    run.font.size = Pt(15)
    run.font.color.rgb = COLOR_WHITE
    run.font.bold = True

    # Arrow + action
    add_text_block(slide, Emu(5000000), row_y + Emu(150000), Emu(11000000), Emu(600000),
                  f"→  {action}", FONT_CONTENT, Pt(16), COLOR_CARD_BODY)

add_text_block(slide, Emu(900000), Emu(9200000), Emu(15000000), Emu(500000),
              "Impact: Shift from 'station calls when tank is empty' to 'delivery arrives before station knows they need it.'",
              FONT_BODY, Pt(16), COLOR_DK1, bold=True)
add_footer(slide)

# Slide: MuleSoft Integration Layer
slide = add_content_slide(prs,
    "MuleSoft: Connecting the Physical and Digital Supply Chain",
    "The unified data bus that connects physical logistics assets to commercial decision-making"
)
integrations = [
    ("TMS", "Route scheduling, fleet allocation, freight costs"),
    ("IoT/Telemetry", "Tank levels, GPS tracking, fuel quality sensors"),
    ("SAP MM/WM", "Inventory at terminals and bases, goods movement"),
    ("ANP Regulatory", "Nota fiscal eletrônica, SEFAZ compliance"),
    ("Third-party Carriers", "EDI integration for contracted fleet scheduling and billing"),
]
for i, (system, desc) in enumerate(integrations):
    x = Emu(900000) + (i % 3) * Emu(5600000)
    y = Emu(2800000) + (i // 3) * Emu(2400000)
    add_card(slide, x, y, Emu(5200000), Emu(2000000), system, desc)
add_footer(slide)

# ============================================================
# SECTION 3: COMMERCIAL EXCELLENCE
# ============================================================

add_section_divider(prs,
    "Growth Through Commercial Excellence",
    "How leading distributors turn disciplined execution into higher volume, stronger retention, and profitable growth."
)

# Slide: Two Battlefields
slide = add_content_slide(prs,
    "Retail (B2C2B) vs. Direct B2B: Two Models, One Platform",
    "The commercial organization fights on two fronts simultaneously"
)
# Retail side
add_text_block(slide, Emu(900000), Emu(2400000), Emu(7500000), Emu(600000),
              "RETAIL — The Branded Network", FONT_CARD, Pt(20), COLOR_DK2, bold=True)
retail_items = [
    "44,000+ stations — network management at continental scale",
    "49% white-flag — constant conversion + retention battle",
    "Dealer relationship: financial support, training, brand compliance",
    "Consumer loyalty programs competing for share of wallet",
    "Convenience retail (AmPm, Shell Select) as margin differentiator"
]
add_multi_text(slide, Emu(900000), Emu(3100000), Emu(7500000), Emu(4000000),
              [f"•  {item}" for item in retail_items], FONT_CARD, Pt(15), COLOR_CARD_BODY)

# B2B side
add_text_block(slide, Emu(9200000), Emu(2400000), Emu(7500000), Emu(600000),
              "DIRECT B2B — The Volume Engine", FONT_CARD, Pt(20), COLOR_DK2, bold=True)
b2b_items = [
    "Agribusiness — seasonal demand, barter contracts, on-farm delivery",
    "Mining & Industry — long-term contracts, dedicated storage (PAVs)",
    "Transport fleets — fleet cards, multi-site, real-time monitoring",
    "Thermal power plants — massive spot volumes during hydro shortage",
    "Government & public — tenders, compliance, payment terms"
]
add_multi_text(slide, Emu(9200000), Emu(3100000), Emu(7500000), Emu(4000000),
              [f"•  {item}" for item in b2b_items], FONT_CARD, Pt(15), COLOR_CARD_BODY)
add_footer(slide)

# Slide: Retail Operating Model
slide = add_content_slide(prs,
    "Retail Network Management: The Five Disciplines",
    "How market leaders manage 10,000+ branded stations at scale"
)
retail_disciplines = [
    ("1. Network Planning", ["Whitespace analysis", "Competitive mapping", "CAPEX prioritization"]),
    ("2. Dealer Acquisition", ["Prospect scoring", "Contract negotiation", "Onboarding workflow"]),
    ("3. Relationship Mgmt", ["Territory segmentation", "Visit cadence & execution", "Business plans per station"]),
    ("4. Volume & Margin", ["Volume forecasting", "Incentive programs", "Product mix optimization"]),
    ("5. Retention & Churn", ["Early warning signals", "Proactive intervention", "Contract renewal"]),
]
for i, (name, items) in enumerate(retail_disciplines):
    x = Emu(900000) + (i * Emu(3400000))
    add_text_block(slide, x, Emu(2600000), Emu(3200000), Emu(600000),
                  name, FONT_CARD, Pt(14), COLOR_HEADER, bold=True)
    item_text = "\n".join([f"• {item}" for item in items])
    add_text_block(slide, x, Emu(3250000), Emu(3200000), Emu(2000000),
                  item_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)
add_footer(slide)

# Slide: B2B Operating Model
slide = add_content_slide(prs,
    "B2B Sales Machine: The Five Disciplines",
    "Complex sales execution for agribusiness, mining, transport, and industrial accounts"
)
b2b_disciplines = [
    ("1. Market Intelligence", ["Sector mapping", "Account identification", "Competitive displacement"]),
    ("2. Complex Sales", ["Multi-stakeholder selling", "Solution design", "Proposal automation"]),
    ("3. Account Development", ["Share-of-wallet growth", "Cross-sell opportunities", "Strategic planning"]),
    ("4. Service Alignment", ["SLA definition & monitoring", "Dedicated infrastructure", "Issue resolution"]),
    ("5. Renewal & Retention", ["Contract lifecycle", "Price renegotiation triggers", "Competitor defense"]),
]
for i, (name, items) in enumerate(b2b_disciplines):
    x = Emu(900000) + (i * Emu(3400000))
    add_text_block(slide, x, Emu(2600000), Emu(3200000), Emu(600000),
                  name, FONT_CARD, Pt(14), COLOR_HEADER, bold=True)
    item_text = "\n".join([f"• {item}" for item in items])
    add_text_block(slide, x, Emu(3250000), Emu(3200000), Emu(2000000),
                  item_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)
add_footer(slide)

# Slide: Commercial Gap
add_gap_slide(prs,
    "The Gap: Two Worlds That Don't Talk to Each Other",
    "Commercial teams operate on instinct without data to scale what the best reps do naturally.",
    [
        ("Customer Data", "Dealer in one system, consumer loyalty in another, B2B in SAP", "🔴 NO 360° VIEW — fragmented customer reality", "red"),
        ("Pipeline Visibility", "Volume forecasts in Excel, opportunities in spreadsheets", "🔴 NO PREDICTABILITY — surprises every quarter", "red"),
        ("Activity Management", "Paper forms, disconnected apps, meeting notes in email", "🟠 NO EXECUTION VISIBILITY — can't scale best practices", "orange"),
        ("Performance Tracking", "Monthly volume reports from SAP, quarterly reviews", "🟠 LATE REACTION — problems found too late", "orange"),
        ("Churn Detection", "Discovered when station unbrands or contract is lost", "🔴 PREVENTABLE LOSSES — no early warning system", "red"),
    ]
)

# Slide: Salesforce Commercial — Retail
add_table_slide(prs,
    "Sales Cloud + Experience Cloud: The Retail Network Platform",
    "Unified platform for managing 44,000+ stations",
    ["Capability", "Salesforce Product", "What It Does"],
    [
        ["Dealer 360° Profile", "Sales Cloud", "Unified view: contract, volume, financials, visits, cases, loyalty performance"],
        ["Territory Intelligence", "Maps + Tableau", "Geo-visualization: network, whitespace, competitor proximity, risk zones"],
        ["Visit Execution", "Sales Cloud Mobile", "Structured visit plans, offline checklists, photo capture, GPS confirmation"],
        ["Dealer Portal", "Experience Cloud", "Self-service: order tracking, financials, marketing materials, support"],
        ["Churn Prediction", "Einstein + Data Cloud", "AI model on historical patterns: volume drop + payment delay + low engagement"],
        ["Incentive Management", "Revenue Cloud", "Automated rebate calculations, tiered programs, transparent dealer dashboards"],
    ]
)

# Slide: Salesforce Commercial — B2B
add_table_slide(prs,
    "Sales Cloud + Revenue Cloud: The B2B Growth Engine",
    "From prospecting to renewal — one connected platform",
    ["Capability", "Salesforce Product", "What It Does"],
    [
        ["Account Planning", "Sales Cloud", "Strategic plans with share-of-wallet, growth targets, stakeholder mapping"],
        ["Pipeline Management", "Sales Cloud", "Stage-gated process with probability-weighted forecasting"],
        ["CPQ & Proposals", "Revenue Cloud", "Dynamic pricing within guardrails, auto-generated indexed contract terms"],
        ["Complex Approvals", "Flow + Agentforce", "Multi-level deal approval with margin validation and audit trail"],
        ["Contract Lifecycle", "Revenue Cloud", "Renewal alerts, price adjustment triggers, amendment workflows"],
        ["Forecast Accuracy", "Einstein Analytics", "AI forecast with consumption patterns, seasonality, macro signals"],
    ]
)

# Slide: Commercial Agentforce
add_agentforce_slide(prs,
    "Commercial Excellence — Agentforce in Action",
    [
        ("Churn Prevention Agent", [
            "Monitors 44,000 stations daily for early warnings",
            "Detects: volume decline >10%, payment delays, reduced mix",
            "Auto-triggers intervention with specific recommended action"
        ]),
        ("Opportunity Spotter Agent", [
            "Scans market: harvest forecasts, mining CAPEX, fleet expansion",
            "Identifies new accounts or expansion opportunities",
            "Creates qualified leads with context and volume estimates"
        ]),
        ("Deal Desk Agent", [
            "Reviews B2B proposals against margin rules and credit limits",
            "Auto-approves standard deals within policy",
            "Escalates exceptions with full context and competitor analysis"
        ]),
    ]
)

# ============================================================
# SECTION 4: CUSTOMER OPERATIONS
# ============================================================

add_section_divider(prs,
    "Loyalty Through Customer Operations Excellence",
    "How distributors create stickiness, differentiate beyond the commodity, and build recurring revenue streams."
)

# Slide: Four Pillars
slide = add_content_slide(prs,
    "Beyond the Commodity: The Four Pillars of Customer Stickiness",
    "Fuel is a commodity. Customer operations transforms it into a relationship business."
)
pillars = [
    ("Credit & Financial Services", "Working capital for dealers, flexible payment for B2B.\nIn a 14% interest rate environment, credit = volume lock-in."),
    ("Loyalty & Rewards", "Consumer programs: Premmia, Shell Box, Km de Vantagens.\n40M+ enrolled consumers = recurring traffic + behavioral data."),
    ("Fleet Management", "Integrated fuel + toll + maintenance cards.\nHigh switching costs for transport companies — guaranteed volume."),
    ("TRR Operations", "Last-mile delivery to farms, small industries, generators.\nServes the fragmented middle-market efficiently."),
]
for i, (name, desc) in enumerate(pillars):
    x = Emu(900000) + (i * Emu(4200000))
    add_card(slide, x, Emu(2600000), Emu(3900000), Emu(3600000), name, desc)
add_footer(slide)

# Slide: Credit Operating Model
slide = add_content_slide(prs,
    "Credit as a Competitive Weapon in a High-Interest Economy",
    "The credit value chain + Brazil-specific complexity"
)
credit_chain = [
    "1. Risk Assessment — Financial analysis, credit scoring, collateral",
    "2. Limit Setting — Dynamic limits based on volume, history, market",
    "3. Disbursement — Working capital lines, product financing, investment packages",
    "4. Monitoring — Real-time exposure tracking, early warning signals",
    "5. Collection — Automated dunning, renegotiation, legal escalation",
]
add_multi_text(slide, Emu(900000), Emu(2600000), Emu(8000000), Emu(4000000),
              credit_chain, FONT_CARD, Pt(16), COLOR_CARD_BODY)

brazil_credit = [
    "Dealer credit — R$1-5M investment packages tied to exclusivity",
    "B2B seasonal — Agribusiness barter (grain-for-fuel) with crop cycle timing",
    "Consigned inventory — Product on dealer premises remains distributor's asset",
    "Receivables anticipation — Factoring card receivables at preferential rates",
]
add_text_block(slide, Emu(9500000), Emu(2600000), Emu(7000000), Emu(600000),
              "Brazil-Specific Complexity:", FONT_CARD, Pt(16), COLOR_HEADER, bold=True)
add_multi_text(slide, Emu(9500000), Emu(3300000), Emu(7000000), Emu(4000000),
              [f"•  {item}" for item in brazil_credit], FONT_CARD, Pt(15), COLOR_CARD_BODY)
add_footer(slide)

# Slide: Loyalty Ecosystem
slide = add_content_slide(prs,
    "Consumer Loyalty in Fuel Retail: Data Is the New Fuel",
    "Major programs in Brazil and the operating model behind them"
)
add_table_slide(prs,
    "Consumer Loyalty in Fuel Retail: Data Is the New Fuel",
    "Major programs in Brazil + the operating model behind them",
    ["Program", "Distributor", "Members", "Key Mechanism"],
    [
        ["Premmia", "Vibra (BR)", "30M+", "Points → marketplace redemption"],
        ["Shell Box", "Raízen", "15M+", "Cashback + digital payments"],
        ["Km de Vantagens", "Ipiranga", "40M+", "Points → airline miles + partners"],
        ["Abastece Aí", "Vibra", "10M+", "Digital payments + cashback"],
    ]
)

# Slide: Fleet Management
slide = add_content_slide(prs,
    "Fleet Cards: The High-Switching-Cost Volume Guarantor",
    "Why fleet management creates defensible lock-in"
)
fleet_includes = [
    "Fuel cards — Branded acceptance network, per-driver controls, real-time authorization",
    "Toll integration — Single card for fuel + tolls (Vale Pedágio obrigatório)",
    "Maintenance — Preventive scheduling tied to mileage/fuel consumption",
    "Analytics — Consumption per vehicle, route efficiency, driver behavior",
    "Compliance — ANTT regulatory compliance, environmental reporting",
]
add_multi_text(slide, Emu(900000), Emu(2600000), Emu(8000000), Emu(4500000),
              [f"•  {item}" for item in fleet_includes], FONT_CARD, Pt(16), COLOR_CARD_BODY)

lock_in = [
    "Integrated into operational systems (TMS, payroll, cost allocation)",
    "Driver rules configured per company (limits/day, per transaction)",
    "Migration = reconfiguring 100s-1000s of drivers + processes",
    "Volume contracts tied to fleet commitment guarantee minimums",
]
add_text_block(slide, Emu(9500000), Emu(2600000), Emu(7000000), Emu(600000),
              "Why It Creates Lock-In:", FONT_CARD, Pt(16), COLOR_HEADER, bold=True)
add_multi_text(slide, Emu(9500000), Emu(3300000), Emu(7000000), Emu(4000000),
              [f"•  {item}" for item in lock_in], FONT_CARD, Pt(15), COLOR_CARD_BODY)
add_footer(slide)

# Slide: Customer Ops Gap
add_gap_slide(prs,
    "The Gap: Four Pillars Operating in Silos",
    "Each pillar generates valuable data that never flows into commercial decisions.",
    [
        ("Credit", "SAP FI for accounting + Excel + email approvals", "🔴 SLOW — credit approval 5-10 days vs. competitor 24h", "red"),
        ("Loyalty", "Separate platform (outsourced) disconnected from CRM", "🔴 NO PERSONALIZATION — generic offers, declining engagement", "red"),
        ("Fleet", "Dedicated system with limited integration to commercial", "🟠 MISSED CROSS-SELL — fleet data unused for account growth", "orange"),
        ("TRR", "Separate operation, managed as cost center with basic tools", "🟠 UNDER-MONETIZED — serving without understanding potential", "orange"),
    ]
)

# Slide: Customer Ops Salesforce Mapping
add_table_slide(prs,
    "The Four Customer Operations Pillars — Salesforce Mapping",
    "One platform connecting credit, loyalty, fleet, and TRR operations",
    ["Pillar", "Salesforce Product", "Key Capability"],
    [
        ["Credit & Financial", "Financial Services Cloud + Flow", "Automated scoring, dynamic limits, approval workflows, portfolio monitoring, early warning alerts"],
        ["Loyalty & Rewards", "Marketing Cloud + Data Cloud + Loyalty Mgmt", "Unified profiles, AI personalization, real-time triggers, multi-channel engagement"],
        ["Fleet Management", "Service Cloud + Experience Cloud + Data Cloud", "Fleet portal, real-time dashboards, predictive maintenance alerts, compliance reporting"],
        ["TRR Operations", "Field Service + Sales Cloud", "Route optimization, customer microsegmentation, upsell recommendations from consumption patterns"],
    ]
)

# Slide: Customer Ops Agentforce
add_agentforce_slide(prs,
    "Customer Operations — Agentforce in Action",
    [
        ("Credit Risk Monitor", [
            "Continuously evaluates portfolio exposure",
            "Detects: Serasa score drop, payment pattern change, volume decline",
            "Recommends: reduce limit, pause orders, trigger collection"
        ]),
        ("Loyalty Engagement Agent", [
            "Identifies members at risk (no visit 30 days, points expiring)",
            "Generates personalized re-engagement offer",
            "Delivers via preferred channel, measures response"
        ]),
        ("Fleet Expansion Agent", [
            "Monitors fleet consumption patterns for growth signals",
            "Detects new drivers, new routes, increased consumption",
            "Auto-generates expansion proposal with estimated volume"
        ]),
    ]
)

# ============================================================
# SECTION 5: UNIFIED ARCHITECTURE
# ============================================================

add_section_divider(prs,
    "The Complete Salesforce Architecture",
    "One Platform, Five Engines, Connected Intelligence"
)

# Slide: Architecture Layers
slide = add_content_slide(prs,
    "One Platform, Five Engines, Connected Intelligence",
    "End-to-end architecture connecting all decision layers"
)
arch_layers = [
    ("AGENTFORCE LAYER", "Supply Agent • Pricing Agent • Logistics Agent • Commercial Agent • Customer Ops Agent", COLOR_DK2),
    ("APPLICATION LAYER", "Sales Cloud • Revenue Cloud • Service Cloud • Field Service • Experience Cloud • Marketing Cloud • Financial Services Cloud • Loyalty Management", COLOR_DK1),
    ("INTELLIGENCE LAYER", "Tableau • CRM Analytics • Einstein AI • Predictive Models", COLOR_HEADER),
    ("DATA LAYER", "Data Cloud — Unified Customer • Supply Signals • Market Data • IoT/Telemetry", COLOR_DK2),
    ("INTEGRATION LAYER", "MuleSoft — SAP ERP • TMS • IoT Platform • ANP/Tax • Petrobras • FX Feeds", COLOR_HEADER),
]
y = Emu(2400000)
for i, (layer, detail, color) in enumerate(arch_layers):
    row_y = y + (i * Emu(1500000))

    # Layer bar
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Emu(900000), row_y, Emu(16000000), Emu(1200000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    # Layer name
    add_text_block(slide, Emu(1200000), row_y + Emu(100000), Emu(4000000), Emu(500000),
                  layer, FONT_BODY, Pt(16), COLOR_WHITE, bold=True)
    # Detail
    add_text_block(slide, Emu(1200000), row_y + Emu(550000), Emu(15000000), Emu(500000),
                  detail, FONT_CONTENT, Pt(13), COLOR_ACCENT1)

add_footer(slide)

# Slide: Roadmap
slide = add_content_slide(prs,
    "Implementation Roadmap: Quick Wins to Full Transformation",
    "Four phases from foundation to autonomous operations"
)
phases = [
    ("Phase 1: Foundation\n(0-6 months)",
     "MuleSoft integration (SAP + core)\nData Cloud deployment\nSales Cloud for retail (dealer 360)",
     "Churn prediction reduces dealer losses 15-20%"),
    ("Phase 2: Commercial\n(6-12 months)",
     "Revenue Cloud (CPQ) for B2B pricing\nExperience Cloud dealer portal\nTableau dashboards",
     "B2B deal cycle: 10 days → 3 days"),
    ("Phase 3: Operations\n(12-18 months)",
     "Field Service for logistics\nLoyalty Management migration\nFinancial Services Cloud for credit",
     "Predictive delivery reduces stockouts 30%"),
    ("Phase 4: Autonomous\n(18-24 months)",
     "Agentforce across all engines\nFull AI decision automation\nEnd-to-end optimization",
     "Agents handle 60% of routine decisions"),
]
for i, (phase, details, quick_win) in enumerate(phases):
    x = Emu(900000) + (i * Emu(4200000))

    # Phase card
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, x, Emu(2600000), Emu(3900000), Emu(6200000)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    # Phase name
    add_text_block(slide, x + Emu(200000), Emu(2800000), Emu(3500000), Emu(900000),
                  phase, FONT_CARD, Pt(15), COLOR_DK2, bold=True)

    # Details
    add_text_block(slide, x + Emu(200000), Emu(3900000), Emu(3500000), Emu(2400000),
                  details, FONT_CARD, Pt(13), COLOR_CARD_BODY)

    # Quick win
    add_text_block(slide, x + Emu(200000), Emu(7000000), Emu(3500000), Emu(600000),
                  "Quick Win:", FONT_CARD, Pt(12), COLOR_HEADER, bold=True)
    add_text_block(slide, x + Emu(200000), Emu(7500000), Emu(3500000), Emu(1000000),
                  quick_win, FONT_CARD, Pt(13), COLOR_DK2)

add_footer(slide)

# Slide: Business Impact
add_table_slide(prs,
    "The Value of Connected Operations",
    "Expected business impact across key performance indicators",
    ["KPI", "Current State", "With Salesforce", "Impact"],
    [
        ["Dealer Churn Rate", "8-12% annually", "4-6% annually", "~R$200M volume retention"],
        ["B2B Deal Cycle", "10-15 days average", "2-3 days average", "4x pipeline velocity"],
        ["Pricing Response Time", "24-48 hours", "< 1 hour", "Margin protection"],
        ["Delivery SLA Compliance", "85-90%", "97%+", "Customer satisfaction"],
        ["Credit Approval Time", "5-10 days", "Same-day", "Win deals first"],
        ["Loyalty Redemption Rate", "15-20%", "35-45%", "Higher engagement"],
        ["Commercial Productivity", "8-10 interactions/week", "20-25/week", "2.5x effectiveness"],
    ]
)

# Slide: Why Salesforce
add_table_slide(prs,
    "Why This Transformation Requires Salesforce",
    "SAP stays as system of record. Salesforce becomes the system of intelligence and engagement.",
    ["Requirement", "Why Not Just ERP?", "Salesforce Advantage"],
    [
        ["Customer 360°", "ERP is transaction-centric", "Built for relationships and lifecycle"],
        ["Speed of Change", "ERP changes are expensive and slow", "Low-code platform: weeks not months"],
        ["AI at Scale", "ERP AI limited to structured data", "Einstein + Agentforce: multi-source, autonomous"],
        ["Channel Unification", "ERP doesn't serve portals or mobile", "Experience Cloud, Mobile, Marketing — native"],
        ["Ecosystem", "ERP is closed", "AppExchange + MuleSoft: 1000s of connectors"],
        ["Time to Value", "ERP projects: 18-36 months", "Phase 1 live in 6 months with measurable ROI"],
    ]
)

# Final quote
add_quote_slide(prs,
    "SAP stays as the system of record. Salesforce becomes the system of intelligence and engagement — the layer where humans and AI make decisions and take action.",
    "Salesforce Platform Vision"
)

# ============================================================
# SAVE
# ============================================================
output_path = "/Users/iaraujo/Downloads/OilGas_Downstream_POV_Expanded_Sections.pptx"
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"  Total slides: {len(prs.slides)}")
