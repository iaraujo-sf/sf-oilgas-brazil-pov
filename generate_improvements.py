"""
Generate Oil&Gas Downstream POV - Improvement Slides
Additional slides to complement v1.1
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# --- CONSTANTS ---
SLIDE_WIDTH = Emu(18288000)
SLIDE_HEIGHT = Emu(10287000)

COLOR_DK1 = RGBColor(0x00, 0x1E, 0x5B)
COLOR_DK2 = RGBColor(0x02, 0x2A, 0xC0)
COLOR_LT2 = RGBColor(0x00, 0xB3, 0xFF)
COLOR_ACCENT1 = RGBColor(0x90, 0xD0, 0xFE)
COLOR_ACCENT2 = RGBColor(0x04, 0xE1, 0xCB)
COLOR_ACCENT6 = RGBColor(0xEA, 0xF5, 0xFE)
COLOR_HEADER = RGBColor(0x12, 0x37, 0x7A)
COLOR_BODY = RGBColor(0x33, 0x33, 0x33)
COLOR_MUTED = RGBColor(0x85, 0x84, 0x81)
COLOR_CARD_BODY = RGBColor(0x38, 0x46, 0x53)
COLOR_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
COLOR_RED = RGBColor(0xCC, 0x00, 0x00)
COLOR_ORANGE = RGBColor(0xE6, 0x8A, 0x00)
COLOR_BG_LIGHT = RGBColor(0xCF, 0xE9, 0xFE)

FONT_TITLE = "Avant Garde Demi SFDC"
FONT_BODY = "Salesforce Sans"
FONT_CARD = "Inter"
FONT_CONTENT = "Montserrat"

MARGIN_LEFT = Emu(900000)


def add_text_block(slide, left, top, width, height, text, font_name=FONT_BODY,
                   font_size=Pt(16), color=COLOR_BODY, bold=False, alignment=PP_ALIGN.LEFT):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = font_size
    run.font.color.rgb = color
    run.font.bold = bold
    return tf


def add_multi_text(slide, left, top, width, height, lines, font_name=FONT_BODY,
                   font_size=Pt(16), color=COLOR_BODY):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        run = p.add_run()
        run.text = line
        run.font.name = font_name
        run.font.size = font_size
        run.font.color.rgb = color
        p.space_after = Pt(8)
    return tf


def add_section_divider(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_DK1

    txBox = slide.shapes.add_textbox(Emu(1800000), Emu(3200000), Emu(14400000), Emu(2400000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(44)
    run.font.color.rgb = COLOR_WHITE

    txBox2 = slide.shapes.add_textbox(Emu(1800000), Emu(5800000), Emu(14400000), Emu(1600000))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = subtitle
    run2.font.name = FONT_BODY
    run2.font.size = Pt(22)
    run2.font.color.rgb = COLOR_ACCENT1
    return slide


def add_content_slide(prs, title, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = COLOR_WHITE

    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Emu(600000), Emu(16000000), Emu(1200000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = title
    run.font.name = FONT_TITLE
    run.font.size = Pt(36)
    run.font.color.rgb = COLOR_DK1

    if subtitle:
        txBox2 = slide.shapes.add_textbox(MARGIN_LEFT, Emu(1400000), Emu(16000000), Emu(800000))
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        run2 = p2.add_run()
        run2.text = subtitle
        run2.font.name = FONT_BODY
        run2.font.size = Pt(20)
        run2.font.color.rgb = COLOR_BODY
    return slide


def add_card(slide, left, top, width, height, header, body, header_color=COLOR_HEADER):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    txBox = slide.shapes.add_textbox(left + Emu(200000), top + Emu(200000),
                                     width - Emu(400000), Emu(600000))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = header
    run.font.name = FONT_CARD
    run.font.size = Pt(16)
    run.font.color.rgb = header_color
    run.font.bold = True

    txBox2 = slide.shapes.add_textbox(left + Emu(200000), top + Emu(900000),
                                      width - Emu(400000), height - Emu(1100000))
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    run2 = p2.add_run()
    run2.text = body
    run2.font.name = FONT_CARD
    run2.font.size = Pt(14)
    run2.font.color.rgb = COLOR_CARD_BODY
    return shape


def add_footer(slide, text="Salesforce POV • Oil&Gas Downstream"):
    txBox = slide.shapes.add_textbox(MARGIN_LEFT, Emu(9600000), Emu(8000000), Emu(400000))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = FONT_CARD
    run.font.size = Pt(12)
    run.font.color.rgb = COLOR_MUTED


# ============================================================
# BUILD PRESENTATION
# ============================================================
prs = Presentation()
prs.slide_width = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# ============================================================
# SLIDE 1: CUSTOMER REFERENCE / COMPARABLE
# ============================================================
slide = add_content_slide(prs,
    "Salesforce in Energy & Downstream: Proven Global Impact",
    "Selected references from fuel distribution, refining, and energy retail"
)

# Reference cards
references = [
    ("Shell", "Global",
     "Unified dealer management, B2B account intelligence, and loyalty integration across 46,000+ stations.\n\n"
     "Impact: Single customer view across retail, fleet, and B2B. Reduced dealer onboarding from weeks to days."),
    ("Repsol", "Europe",
     "End-to-end commercial transformation: pricing governance, B2B pipeline, loyalty (Waylet) and service operations.\n\n"
     "Impact: 360° customer view, 30% improvement in campaign conversion, unified dealer/B2B platform."),
    ("Cepsa", "Spain/Portugal",
     "Digital transformation of fuel retail network, loyalty program modernization, and fleet card management.\n\n"
     "Impact: 2x loyalty engagement, reduced churn in branded network, data-driven station investment."),
    ("bp / Castrol", "Global",
     "B2B sales transformation for lubricants and fuels. Account planning, CPQ, and predictive service.\n\n"
     "Impact: 25% pipeline growth, 40% faster quote-to-close, unified global account intelligence."),
]

for i, (company, region, description) in enumerate(references):
    col = i % 2
    row = i // 2
    x = Emu(900000) + col * Emu(8400000)
    y = Emu(2600000) + row * Emu(3600000)

    # Card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(7800000), Emu(3200000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    # Company + Region
    add_text_block(slide, x + Emu(250000), y + Emu(200000), Emu(5000000), Emu(500000),
                  company, FONT_CARD, Pt(18), COLOR_DK2, bold=True)
    add_text_block(slide, x + Emu(5500000), y + Emu(250000), Emu(2000000), Emu(400000),
                  region, FONT_CARD, Pt(13), COLOR_MUTED, bold=False)

    # Description
    add_text_block(slide, x + Emu(250000), y + Emu(800000), Emu(7300000), Emu(2200000),
                  description, FONT_CARD, Pt(14), COLOR_CARD_BODY)

add_text_block(slide, MARGIN_LEFT, Emu(9200000), Emu(16000000), Emu(400000),
              "Source: Publicly available Salesforce customer stories and industry references. Specific metrics are illustrative of published outcomes.",
              FONT_CARD, Pt(11), COLOR_MUTED)


# ============================================================
# SLIDE 2: CUSTOMER OPS — OPERATING GAP (was missing)
# ============================================================
slide = add_content_slide(prs,
    "Customer Operations: The Gap That Erodes Stickiness",
    "Four pillars generating data that never becomes coordinated action"
)

gap_items = [
    ("1 · Credit Operations",
     "SAP FI for accounting + Excel for analysis + email approval chains",
     "🔴 SLOW DECISIONS — 5-10 day approval vs. competitor 24h response",
     COLOR_RED),
    ("2 · Loyalty Programs",
     "Outsourced platform disconnected from CRM, commercial and service data",
     "🔴 NO PERSONALIZATION — generic offers, declining engagement, invisible ROI",
     COLOR_RED),
    ("3 · Fleet Management",
     "Dedicated card system with limited integration to sales or account teams",
     "🟠 MISSED CROSS-SELL — fleet consumption data unused for growth plays",
     COLOR_ORANGE),
    ("4 · TRR / Last-Mile",
     "Separate operation managed as cost center with basic route tools",
     "🟠 UNDER-MONETIZED — serving customers without understanding their full value",
     COLOR_ORANGE),
]

y_start = Emu(2600000)
for i, (layer_name, reality, risk_text, risk_color) in enumerate(gap_items):
    y = y_start + (i * Emu(1600000))

    add_text_block(slide, Emu(900000), y, Emu(5000000), Emu(500000),
                  layer_name, FONT_BODY, Pt(22), COLOR_DK1, bold=True)
    add_text_block(slide, Emu(900000), y + Emu(550000), Emu(8000000), Emu(400000),
                  reality, FONT_BODY, Pt(17), COLOR_BODY)
    add_text_block(slide, Emu(11500000), y + Emu(200000), Emu(6000000), Emu(700000),
                  risk_text, FONT_BODY, Pt(14), risk_color, bold=True)

add_text_block(slide, MARGIN_LEFT, Emu(9100000), Emu(16000000), Emu(500000),
              "The core problem: each pillar generates valuable behavioral data about customers — but that data never flows into commercial decisions, proactive service, or retention actions.",
              FONT_BODY, Pt(16), COLOR_DK1, bold=True)
add_footer(slide)


# ============================================================
# SLIDE 3: CUSTOMER OPS — AGENTFORCE
# ============================================================
slide = add_content_slide(prs,
    "Customer Operations — Agentforce in Action",
    "Autonomous agents that turn credit, loyalty, fleet and service signals into governed retention and growth actions"
)

agents = [
    ("Credit Risk Monitor Agent",
     [
         "Continuously evaluates portfolio exposure across dealers, fleets and B2B",
         "Detects: Serasa score change, payment pattern shift, volume decline + high balance",
         "Recommends: limit adjustment, order pause, collection trigger, or proactive renegotiation",
         "Documents rationale, risk level, account context and approval path"
     ]),
    ("Loyalty Engagement Agent",
     [
         "Identifies members at risk: no visit in 30 days, points expiring, declining frequency",
         "Generates personalized re-engagement offer based on historical behavior and value tier",
         "Delivers via preferred channel (push, SMS, email, pump screen)",
         "Measures response, adjusts strategy, feeds learning back to segmentation"
     ]),
    ("Fleet Expansion Agent",
     [
         "Monitors fleet customer consumption for growth signals (new drivers, new routes, +volume)",
         "Detects fleet growth above threshold or contract utilization approaching limit",
         "Alerts account manager with expansion recommendation + estimated additional volume",
         "Auto-generates contract amendment proposal with updated pricing corridor"
     ]),
    ("TRR Graduation Agent",
     [
         "Identifies TRR customers with consistent growth above commercial threshold",
         "Evaluates: volume trajectory, payment history, location, credit profile, product mix",
         "Recommends 'graduation' to direct commercial relationship with full context",
         "Triggers prospecting workflow with consumption history, margin analysis and credit pre-score"
     ]),
]

for i, (agent_name, bullets) in enumerate(agents):
    col = i % 2
    row = i // 2
    x = Emu(900000) + col * Emu(8400000)
    y = Emu(2600000) + row * Emu(3600000)

    # Card bg
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(7800000), Emu(3300000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    # Agent name
    add_text_block(slide, x + Emu(250000), y + Emu(200000), Emu(7300000), Emu(500000),
                  agent_name, FONT_CARD, Pt(16), COLOR_DK2, bold=True)

    # Bullets
    bullet_text = "\n".join([f"•  {b}" for b in bullets])
    add_text_block(slide, x + Emu(250000), y + Emu(800000), Emu(7300000), Emu(2300000),
                  bullet_text, FONT_CARD, Pt(13), COLOR_CARD_BODY)

add_text_block(slide, MARGIN_LEFT, Emu(9400000), Emu(16000000), Emu(400000),
              "Operating principle: every agent recommendation links trigger, customer context, policy, owner, action and KPI outcome inside the Salesforce workflow.",
              FONT_BODY, Pt(14), COLOR_MUTED)


# ============================================================
# SLIDE 4: ENERGY & UTILITIES CLOUD
# ============================================================
slide = add_content_slide(prs,
    "Energy & Utilities Cloud: The Industry Accelerator",
    "Why a pre-built industry data model matters for time-to-value in downstream"
)

# Left side - What E&U Cloud provides
add_text_block(slide, MARGIN_LEFT, Emu(2400000), Emu(8000000), Emu(600000),
              "What Energy & Utilities Cloud provides out-of-the-box:", FONT_CARD, Pt(16), COLOR_HEADER, bold=True)

eu_features = [
    "Industry data model: Accounts, Sites, Assets, Service Points, Contracts hierarchy",
    "Account-to-Meter / Account-to-Location relationships native to energy",
    "Pre-built processes: service agreements, usage tracking, regulatory compliance",
    "Asset lifecycle management for tanks, terminals, and delivery infrastructure",
    "Multi-site B2B account structures (mining, agribusiness, fleet companies)",
    "Regulatory and compliance workflows (ANP, environmental, CBIOs)"
]
add_multi_text(slide, MARGIN_LEFT, Emu(3100000), Emu(8000000), Emu(4500000),
              [f"•  {f}" for f in eu_features], FONT_CARD, Pt(15), COLOR_CARD_BODY)

# Right side - Why it matters
add_text_block(slide, Emu(9500000), Emu(2400000), Emu(7500000), Emu(600000),
              "Why this matters vs. generic Sales/Service Cloud:", FONT_CARD, Pt(16), COLOR_HEADER, bold=True)

comparisons = [
    ("Generic CRM", "6-9 months to build industry data model from scratch"),
    ("E&U Cloud", "Industry model pre-built — focus on business logic, not schema"),
    ("Generic CRM", "Custom objects for sites, assets, service points"),
    ("E&U Cloud", "Native relationships: Account → Site → Asset → Service Agreement"),
    ("Generic CRM", "Manual integration design for metering/telemetry"),
    ("E&U Cloud", "Pre-built patterns for usage data, consumption, and IoT signals"),
]

y = Emu(3100000)
for i, (label, desc) in enumerate(comparisons):
    row_y = y + (i * Emu(700000))
    color = COLOR_MUTED if label == "Generic CRM" else COLOR_DK2
    style = False if label == "Generic CRM" else True
    prefix = "✗  " if label == "Generic CRM" else "✓  "
    add_text_block(slide, Emu(9500000), row_y, Emu(7500000), Emu(500000),
                  f"{prefix}{desc}", FONT_CARD, Pt(14), color, bold=style)

# Bottom message
add_text_block(slide, MARGIN_LEFT, Emu(8600000), Emu(16000000), Emu(600000),
              "Bottom line: Energy & Utilities Cloud reduces implementation time by 30-40% for the first wave by providing industry-native data model, processes, and integration patterns.",
              FONT_BODY, Pt(16), COLOR_DK1, bold=True)
add_footer(slide)


# ============================================================
# SLIDE 5: INTEGRATION LANDSCAPE
# ============================================================
slide = add_content_slide(prs,
    "Integration Landscape: Connecting the Downstream Ecosystem",
    "MuleSoft as the API & event fabric connecting 8+ system categories to the Salesforce intelligence layer"
)

# System categories with specific platforms
systems = [
    ("ERP / Core Finance", "SAP S/4HANA, SAP IS-Oil&Gas\nSAP MM, FI, CO-PA, SD",
     "Master data, costs, invoicing, credit, inventory postings"),
    ("Transport Management", "SAP TM, Logix, TOTVS Logística\nCustom TMS platforms",
     "Routes, schedules, freight costs, carrier allocation, delivery status"),
    ("Pricing & Market", "Petrobras price feeds, ANP\nB3 (FX), Bloomberg, Argus",
     "Gate prices, FX rates, commodity benchmarks, competitor monitoring"),
    ("IoT / Telemetry", "Tank level sensors, GPS fleet\nSCADA, flow meters",
     "Real-time tank levels, vehicle tracking, quality measurement"),
    ("Fiscal & Tax", "SEFAZ (NF-e), SPED\nICMS engines, compliance",
     "Tax documents, fiscal compliance, state-specific obligations"),
    ("Loyalty & Payments", "Card processors, acquirers\nLoyalty platforms (legacy)",
     "Transaction data, points, redemptions, fleet card authorizations"),
    ("Regulatory", "ANP systems, IBAMA\nCBIOs / RenovaBio platform",
     "Compliance reporting, blend certificates, environmental credits"),
    ("Third-party / Carriers", "EDI with transporters\nPort/terminal operators",
     "Availability, capacity, scheduling, proof of delivery, incidents"),
]

for i, (category, platforms, data_flow) in enumerate(systems):
    col = i % 4
    row = i // 4
    x = Emu(900000) + col * Emu(4200000)
    y = Emu(2600000) + row * Emu(3400000)

    # Card
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, Emu(3900000), Emu(3000000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_ACCENT6
    shape.line.fill.background()

    # Category header
    add_text_block(slide, x + Emu(200000), y + Emu(150000), Emu(3500000), Emu(500000),
                  category, FONT_CARD, Pt(14), COLOR_DK2, bold=True)
    # Platforms
    add_text_block(slide, x + Emu(200000), y + Emu(700000), Emu(3500000), Emu(900000),
                  platforms, FONT_CARD, Pt(12), COLOR_CARD_BODY)
    # Data flow
    add_text_block(slide, x + Emu(200000), y + Emu(1800000), Emu(3500000), Emu(1000000),
                  f"→ {data_flow}", FONT_CARD, Pt(12), COLOR_MUTED)

# MuleSoft bar at bottom
shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,
                               Emu(900000), Emu(9000000), Emu(16400000), Emu(700000))
shape.fill.solid()
shape.fill.fore_color.rgb = COLOR_DK2
shape.line.fill.background()
add_text_block(slide, Emu(1200000), Emu(9150000), Emu(15800000), Emu(400000),
              "MuleSoft Anypoint Platform — Reusable APIs • Event-Driven Architecture • Canonical Data Services • API Governance • Composable Integration",
              FONT_BODY, Pt(15), COLOR_WHITE, bold=True, alignment=PP_ALIGN.CENTER)


# ============================================================
# SLIDE 6: INVESTMENT FRAMEWORK
# ============================================================
slide = add_content_slide(prs,
    "Investment Framework: T-Shirt Sizing by Wave",
    "Indicative investment ranges to support internal budget positioning"
)

# Table-like layout
headers = ["Wave", "Scope", "Duration", "Team Size", "Investment Range", "Expected ROI Trigger"]
rows = [
    ["Wave 1\nFoundation", "MuleSoft core APIs + Data Cloud\nSales Cloud (Dealer 360)\nTableau executive dashboards",
     "4-6\nmonths", "8-12\nresources", "$$\n(R$2-4M)", "Churn prediction\noperational within\n90 days of go-live"],
    ["Wave 2\nCommercial", "Revenue Cloud (CPQ/Pricing)\nExperience Cloud (Dealer Portal)\nB2B pipeline governance",
     "5-7\nmonths", "10-15\nresources", "$$$\n(R$4-7M)", "B2B deal cycle\nreduction visible\nwithin 60 days"],
    ["Wave 3\nOperations", "Field Service + Logistics\nFinancial Services Cloud (Credit)\nLoyalty Management migration",
     "6-8\nmonths", "12-18\nresources", "$$$\n(R$5-8M)", "Delivery SLA\nimprovement within\nfirst quarter"],
    ["Wave 4\nAutonomous", "Agentforce deployment\nAI-driven decision automation\nFull operating model scale",
     "4-6\nmonths", "8-12\nresources", "$$\n(R$3-5M)", "60% routine decisions\nautomated within\n6 months"],
]

# Header row
y_header = Emu(2400000)
col_widths = [Emu(2200000), Emu(5000000), Emu(1600000), Emu(1600000), Emu(2400000), Emu(3600000)]
col_starts = [Emu(900000)]
for w in col_widths[:-1]:
    col_starts.append(col_starts[-1] + w)

for i, (header, x, w) in enumerate(zip(headers, col_starts, col_widths)):
    add_text_block(slide, x, y_header, w, Emu(500000),
                  header, FONT_CARD, Pt(13), COLOR_WHITE, bold=True)
    # Header bg
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_header, w - Emu(50000), Emu(500000))
    shape.fill.solid()
    shape.fill.fore_color.rgb = COLOR_DK1
    shape.line.fill.background()
    # Re-add text on top
    add_text_block(slide, x + Emu(100000), y_header + Emu(100000), w - Emu(200000), Emu(400000),
                  header, FONT_CARD, Pt(12), COLOR_WHITE, bold=True)

# Data rows
for r, row in enumerate(rows):
    y_row = Emu(3100000) + r * Emu(1600000)
    bg_color = COLOR_ACCENT6 if r % 2 == 0 else COLOR_WHITE

    for c, (cell, x, w) in enumerate(zip(row, col_starts, col_widths)):
        # Background
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y_row, w - Emu(50000), Emu(1400000))
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg_color
        shape.line.fill.background()

        font_size = Pt(11) if c > 0 else Pt(12)
        bold = True if c == 0 else False
        color = COLOR_DK2 if c == 0 else COLOR_CARD_BODY
        add_text_block(slide, x + Emu(100000), y_row + Emu(100000),
                      w - Emu(200000), Emu(1200000),
                      cell, FONT_CARD, font_size, color, bold)

# Disclaimer
add_text_block(slide, MARGIN_LEFT, Emu(9400000), Emu(16000000), Emu(400000),
              "Note: Ranges are indicative and depend on system complexity, data readiness, and organizational change management. Formal scoping follows Discovery & Design phase.",
              FONT_CARD, Pt(11), COLOR_MUTED)


# ============================================================
# SLIDE 7: REFINED AGENDA (suggestion)
# ============================================================
slide = add_content_slide(prs,
    "Suggested Agenda Refinement",
    "More assertive titles that frame each section as a business decision, not just information"
)

agenda_items = [
    ("01", "Current Agenda", "Brazil Downstream Market Dynamics",
     "Suggested", "A $76B Market Where Decision Speed Defines Margin"),
    ("02", "Current Agenda", "The Five Downstream Value Engines",
     "Suggested", "Five Engines That Create (or Leak) Value Every Day"),
    ("03", "Current Agenda", "Operating Gaps and Decision Challenges",
     "Suggested", "The Hidden Cost of Disconnected Operations"),
    ("04", "Current Agenda", "Salesforce as the Intelligence & Engagement Layer",
     "Suggested", "From System of Record to System of Intelligent Action"),
    ("05", "Current Agenda", "Business Engines Enabled by Salesforce",
     "Suggested", "Connected Excellence: Pricing, Logistics, Commercial & Customer Ops"),
    ("06", "Current Agenda", "Transformation Roadmap and Value Realization",
     "Suggested", "Quick Wins to Full Transformation: A Governed Path"),
]

y = Emu(2400000)
for i, (num, label_curr, current, label_new, suggested) in enumerate(agenda_items):
    row_y = y + i * Emu(1100000)

    # Number
    add_text_block(slide, Emu(900000), row_y, Emu(600000), Emu(500000),
                  num, FONT_TITLE, Pt(20), COLOR_DK2, bold=True)

    # Current (strikethrough style - gray)
    add_text_block(slide, Emu(1500000), row_y, Emu(7000000), Emu(500000),
                  current, FONT_CARD, Pt(14), COLOR_MUTED)

    # Arrow
    add_text_block(slide, Emu(8800000), row_y + Emu(50000), Emu(600000), Emu(400000),
                  "→", FONT_BODY, Pt(18), COLOR_DK2, bold=True)

    # Suggested (bold, blue)
    add_text_block(slide, Emu(9400000), row_y, Emu(7500000), Emu(500000),
                  suggested, FONT_CARD, Pt(14), COLOR_DK2, bold=True)

add_text_block(slide, MARGIN_LEFT, Emu(9200000), Emu(16000000), Emu(400000),
              "Principle: each agenda title should frame a business decision or tension, not just describe content. The audience should feel 'I need to hear this.'",
              FONT_BODY, Pt(14), COLOR_MUTED)


# ============================================================
# SAVE
# ============================================================
output_path = "/Users/iaraujo/Downloads/OilGas_Downstream_POV_Improvements.pptx"
prs.save(output_path)
print(f"✓ Presentation saved: {output_path}")
print(f"  Total slides: {len(prs.slides)}")
